"""Add an 'Introduction · VSR Applications' slide before the technical Motivation.

The slide leaves empty rounded-rectangle placeholders where the author will
paste her real LR/HR image pairs. Each pair has a colored arrow between
LR and HR, plus a category label and one-line caption.

Layout — 2×2 grid:

   Surveillance          Remote sensing
   [LR] → [HR]           [LR] → [HR]

   Medical imaging       Autonomous driving
   [LR] → [HR]           [LR] → [HR]

After insert, the slide is moved to position 3 (0-indexed 2) so the
chronological deck flow becomes:
  1. Cover · 2. Contents · 3. Applications (NEW) · 4. Motivation · …
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SRC = Path("/home/user/StableVSR/presentation/WC_BD_SFT_Defense_Cho_YuJin_EN.pptx")


# Palette — matches the EN deck
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x12, 0x1E, 0x3D)
MUTED_INK = RGBColor(0x3D, 0x3D, 0x4A)
GRAY = RGBColor(0x6B, 0x6B, 0x78)
LIGHT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
PALE = RGBColor(0xF4, 0xF4, 0xF8)
DEEP = RGBColor(0x15, 0x12, 0x7C)

# Section 1 accent colors for the 4 categories (same hue family,
# different shades — matches the typical "intro slide" reference image).
COL_RED   = RGBColor(0xCB, 0x42, 0x35)    # Surveillance
COL_BLUE  = RGBColor(0x2C, 0x55, 0xA8)    # Remote sensing
COL_GREEN = RGBColor(0x1F, 0x6E, 0x5C)    # Medical imaging
COL_GOLD  = RGBColor(0xC5, 0x8A, 0x12)    # Autonomous driving


def _rm(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def add_textbox(slide, left, top, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    return tb


def add_para(tf, text, *, size=13, bold=False, italic=False, color=INK,
              align=PP_ALIGN.LEFT, space_before=2, space_after=4):
    if not tf.text and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    if color is not None:
        r.font.color.rgb = color
    return p


def add_placeholder_rect(slide, left, top, w, h, *, label="LR"):
    """Empty rounded rectangle with dashed border + faint label.
    User will paste an image over this; the rectangle gets covered."""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(w), Inches(h),
    )
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = PALE
    sh.line.color.rgb = GRAY
    sh.line.width = Pt(0.75)
    # Faint placeholder label
    sh.text_frame.text = ""
    p = sh.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = p.add_run()
    r.text = f"[ paste {label} image ]"
    r.font.size = Pt(10)
    r.font.italic = True
    r.font.color.rgb = GRAY
    return sh


def add_arrow(slide, x0, y0, x1, y1, *, color):
    """Solid colored right-pointing arrow connector between two LR/HR boxes."""
    sh = slide.shapes.add_connector(
        2,  # MSO_CONNECTOR.STRAIGHT
        Inches(x0), Inches(y0), Inches(x1), Inches(y1),
    )
    sh.line.color.rgb = color
    sh.line.width = Pt(3.5)
    # End arrow head
    from pptx.oxml.ns import qn
    ln = sh.line._get_or_add_ln()
    headEnd = ln.find(qn('a:headEnd'))
    tailEnd = ln.find(qn('a:tailEnd'))
    if tailEnd is None:
        from lxml import etree
        tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
    tailEnd.set('type', 'triangle')
    tailEnd.set('w', 'med')
    tailEnd.set('len', 'med')


def add_label(slide, left, top, w, h, text, color):
    """Category label beneath an LR/HR pair."""
    tb = add_textbox(slide, left, top, w, h, anchor=MSO_ANCHOR.MIDDLE)
    add_para(tb.text_frame, text, size=12, bold=True, color=color,
              align=PP_ALIGN.CENTER, space_before=0, space_after=0)


def build_apps_slide(prs):
    """Build the Applications slide using layout 12 ('16_제목만') for the
    deck's house style. The user can move/duplicate it in PowerPoint."""
    layout = prs.slide_masters[0].slide_layouts[12]
    s = prs.slides.add_slide(layout)
    # Remove the layout's default body placeholder
    for ph in list(s.placeholders):
        if ph.placeholder_format.idx == 10:
            _rm(ph)

    # Title (white text on layout's navy bar)
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = "Introduction"
            r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
            r2 = p.add_run()
            r2.text = "   ·   Where VSR is used"
            r2.font.size = Pt(14); r2.font.color.rgb = RGBColor(0xDD, 0xD8, 0xF0)
            break

    # Sub-header — small "❑ Motivation" style marker (mirrors author's reference)
    tb = add_textbox(s, 0.55, 0.95, 12.3, 0.45)
    add_para(tb.text_frame, "❑   Motivation",
              size=14, bold=True, color=INK)

    # Brief one-liner caption
    tb = add_textbox(s, 0.55, 1.40, 12.3, 0.4)
    add_para(tb.text_frame,
              "VSR enables a wide range of real-world applications where "
              "low-resolution input must be enhanced for downstream tasks.",
              size=11, italic=True, color=GRAY)

    # ── 2×2 grid of application examples ───────────────────────────
    # Cell geometry: each cell holds [LR rect] → [HR rect] + label.
    pair_box_w = 1.9      # LR / HR rectangle width
    pair_box_h = 1.55     # rectangle height
    arrow_gap_w = 0.55    # space between LR and HR for the arrow

    # Cell origins (top-left corner of each LR rect)
    rows_y = [2.05, 4.60]
    cols_x = [0.95, 7.55]
    label_y_off = pair_box_h + 0.08

    apps = [
        ("Surveillance · Face / License plate",
         "low-quality CCTV → identifiable",
         COL_RED),
        ("Remote sensing · Satellite imagery",
         "wide-area low-res → asset detail",
         COL_BLUE),
        ("Medical imaging · MRI / CT",
         "fast acquisition → fine structure",
         COL_GREEN),
        ("Autonomous driving · Traffic-sign / pedestrian",
         "distant low-res → reliable detection",
         COL_GOLD),
    ]

    for i, (label_main, label_sub, color) in enumerate(apps):
        col = i % 2
        row = i // 2
        x = cols_x[col]
        y = rows_y[row]

        # LR placeholder
        add_placeholder_rect(s, x, y, pair_box_w, pair_box_h, label="LR")
        # Arrow
        arrow_x0 = x + pair_box_w + 0.05
        arrow_x1 = x + pair_box_w + arrow_gap_w - 0.05
        arrow_y = y + pair_box_h / 2
        add_arrow(s, arrow_x0, arrow_y, arrow_x1, arrow_y, color=color)
        # HR placeholder
        hr_x = x + pair_box_w + arrow_gap_w
        add_placeholder_rect(s, hr_x, y, pair_box_w, pair_box_h, label="HR")

        # Category label below the pair
        cell_total_w = pair_box_w * 2 + arrow_gap_w
        add_label(s, x, y + label_y_off, cell_total_w, 0.35,
                   label_main, color=color)
        # Sub-caption
        tb = add_textbox(s, x, y + label_y_off + 0.35, cell_total_w, 0.30)
        add_para(tb.text_frame, label_sub,
                  size=10, italic=True, color=GRAY,
                  align=PP_ALIGN.CENTER, space_before=0, space_after=0)

    # Bottom takeaway / hand-off to the next slide
    bottom_y = 7.00
    tb = add_textbox(s, 0.55, bottom_y, 12.3, 0.30)
    add_para(tb.text_frame,
              "→  Across all these domains, video VSR adds a uniquely hard "
              "constraint: temporal consistency across frames.",
              size=10.5, italic=True, color=INK, align=PP_ALIGN.CENTER)

    # ── Korean speaker notes ───────────────────────────────────────
    ntf = s.notes_slide.notes_text_frame
    ntf.clear()
    p = ntf.paragraphs[0]
    p.add_run().text = (
        "Introduction · VSR이 어디에 쓰이는가 — 실제 application 예시 "
        "슬라이드입니다.\n\n"
        "네 가지 대표 영역으로 VSR이 사용됩니다.\n\n"
        "첫째 — Surveillance. CCTV 같은 저화질 입력에서 얼굴 / 번호판 "
        "같은 식별 가능한 detail을 복원하는 forensic / security 응용.\n\n"
        "둘째 — Remote sensing. 위성 / 항공 영상의 광범위한 저해상도 "
        "촬영에서 구조물 / 자산 단위의 detail을 분석 가능한 수준으로 "
        "끌어올리는 응용.\n\n"
        "셋째 — Medical imaging. MRI / CT 등에서 빠른 촬영을 위해 "
        "낮은 해상도로 획득한 영상을 fine structure 분석 가능한 "
        "수준으로 복원하는 응용. 의사의 진단 보조에 활용.\n\n"
        "넷째 — Autonomous driving. 멀리 떨어진 traffic sign, 보행자 "
        "같은 작은 객체를 detection 가능한 해상도로 복원해 인식 "
        "안정성을 높이는 응용.\n\n"
        "이런 영상 기반 응용에서는 single image SR과 달리 video VSR "
        "이 가지는 고유한 어려움이 있습니다 — temporal consistency. "
        "다음 슬라이드에서 그 문제를 자세히 다루겠습니다."
    )

    return s


def reorder_to_position(prs, new_pos_0idx):
    """Move the most-recently-appended slide to 0-indexed position."""
    sldIdLst = prs.slides._sldIdLst
    new_sid = sldIdLst[-1]
    sldIdLst.remove(new_sid)
    sldIdLst.insert(new_pos_0idx, new_sid)


def main():
    prs = Presentation(str(SRC))
    print(f"Before: {len(prs.slides)} slides")
    build_apps_slide(prs)
    # Place the new slide AFTER Contents (0-indexed 1), so it becomes
    # slide 3 in the final deck.
    reorder_to_position(prs, 2)
    prs.save(str(SRC))
    print(f"After:  {len(prs.slides)} slides")
    print(f"Saved (overwritten): {SRC}")


if __name__ == "__main__":
    main()
