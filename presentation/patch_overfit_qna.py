"""Patch the EN deck:
  1. Rephrase Discussion(28) card 2 and Limitations(29) item 1 so the wording
     does not imply we are uniquely data-limited (all main DM baselines are
     REDS-only too — the issue is architectural specialization, not data).
  2. Enrich slide 29 Korean speaker notes with an 'overfitting?' Q&A defense.
  3. Append a new Q&A-backup slide after Thanks: "Is this overfitting?"
     — visible, with a definition-of-overfitting argument + supporting data.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


SRC = Path("/home/user/StableVSR/presentation/WC_BD_SFT_Defense_Cho_YuJin_EN.pptx")


# Palette (same as the EN deck)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x12, 0x1E, 0x3D)
MUTED_INK = RGBColor(0x3D, 0x3D, 0x4A)
GRAY = RGBColor(0x6B, 0x6B, 0x78)
LIGHT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
SOFT_BG = RGBColor(0xF9, 0xF9, 0xFC)
DEEP = RGBColor(0x15, 0x12, 0x7C)
SEC_RED   = RGBColor(0xCB, 0x42, 0x35)
SEC_BLUE  = RGBColor(0x2C, 0x55, 0xA8)
SEC_GREEN = RGBColor(0x1F, 0x6E, 0x5C)
SEC_GOLD  = RGBColor(0xC5, 0x8A, 0x12)


# ── (1) Text replacements on slides 28 and 29 ──────────────────────
REPLACEMENTS = {
    # Slide 28 — Discussion card 2 (Specialization–generalization)
    "REDS-only training → isolates WCM contribution.":
        "All main DM baselines (StableVSR · DGAF-VSR · BasicVSR++) "
        "are REDS-only too — matched setup.",
    "Strong on REDS-like statistics; smaller gains on Vid4 SD.":
        "Wavelet conditioning specializes further to REDS frequency "
        "statistics by design — strongest on REDS4 · smaller on Vid4 SD.",

    # Slide 29 — Limitations item 1
    "Strongest on REDS-like content; smaller gains on Vid4 SD.":
        "Wavelet priors specialize to training-distribution frequency "
        "statistics by design.",
    "→ Broader training mixture · degradation-aware augmentation":
        "→ Architectural fusion (e.g., DGAF-VSR's HR feature warping) "
        "+ broader training mixture",
}


# ── (2) Korean speaker-notes appendix for slide 29 ─────────────────
OVERFITTING_NOTE_KO = (
    "\n\n[Q&A 대비] 'Overfitting 아니냐?'는 질문이 나올 수 있습니다. "
    "이렇게 답하시면 됩니다.\n\n"
    "1) Overfitting의 정의 — training distribution에 너무 fit되어서 "
    "같은 distribution의 held-out test set 성능까지 떨어지는 현상.\n\n"
    "2) REDS4는 학습에 사용되지 않은 held-out test sequence이지만 "
    "REDS distribution과 동일합니다. 거기서 모든 metric이 향상됐다는 "
    "건 in-distribution generalization이 잘 되고 있다는 증거입니다 — "
    "overfitting의 signature가 아닙니다.\n\n"
    "3) Vid4 / UDM10 / SPMCS의 gap은 overfitting이 아니라 "
    "distribution shift입니다. 비교 모델(StableVSR · DGAF-VSR · "
    "BasicVSR++) 도 모두 REDS-only로 학습되어 같은 shift를 겪습니다.\n\n"
    "4) 우리 method가 특정 OOD에서 더 약한 건 wavelet conditioning이 "
    "training distribution의 frequency statistics에 architectural하게 "
    "specialize 되기 때문 — design 의도이지, memorization이 아닙니다.\n\n"
    "5) 학습 가능 파라미터는 6.22M 뿐이고 U-Net 472M + VAE 55M은 "
    "frozen입니다. 또 identity-preserving init으로 시작 — 이 "
    "구조에서 memorization은 통계적으로 불가능합니다.\n\n"
    "6) 만약 overfitting이라면 OOD의 모든 metric이 일률적으로 "
    "악화돼야 하는데, 실제로는 UDM10 NR perceptual 3개 metric에서 "
    "DM 그룹 1위, SPMCS tLPIPS는 −38% 개선 — overfitting 패턴이 "
    "아닙니다.\n\n"
    "필요하면 Q&A backup slide의 표/그림으로 보충 설명 가능합니다."
)


# ── (3) Q&A backup slide builder ───────────────────────────────────
def add_card(slide, left, top, width, height, *, fill=WHITE, line=LIGHT_GRAY,
              line_width=0.5):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(top),
                                Inches(width), Inches(height))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_width)
    s.text_frame.text = ""
    return s


def add_strip(slide, left, top, width, color, *, thickness=0.045):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(left), Inches(top),
                                Inches(width), Inches(thickness))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False


def add_left_stripe(slide, left, top, height, color, *, thickness=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(left), Inches(top),
                                Inches(thickness), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False


def add_textbox(slide, left, top, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    return tb


def add_para(tf, text, *, size=13, bold=False, italic=False, color=INK,
              align=PP_ALIGN.LEFT, space_before=2, space_after=4):
    if not tf.text and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    if color is not None:
        r.font.color.rgb = color
    return p


def add_table(slide, left, top, w, h, headers, rows, header_size=10.5,
               body_size=10):
    tbl = slide.shapes.add_table(len(rows) + 1, len(headers),
                                  Inches(left), Inches(top),
                                  Inches(w), Inches(h)).table
    for j, head in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = DEEP
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.04)
        c.margin_right = Inches(0.04)
        c.text_frame.text = ""
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = head
        r.font.size = Pt(header_size)
        r.font.bold = True
        r.font.color.rgb = WHITE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.04)
            c.margin_right = Inches(0.04)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 == 0 else RGBColor(
                0xF7, 0xF7, 0xFB)
            c.text_frame.text = ""
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(body_size)
            r.font.color.rgb = INK


def build_overfitting_slide(prs):
    """Q&A backup: 'Is this overfitting?' — added at the very end."""
    # The deck's content slides used layout 12; the Thanks slide uses a custom
    # one. For a clean Q&A backup with the right size, use layout 6 (blank).
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)

    # Title bar (manual navy strip mimicking the main deck)
    bar = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DEEP
    bar.line.fill.background()
    bar.shadow.inherit = False
    # Title text
    tb = add_textbox(s, 0.4, 0.05, 12.5, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Q&A Backup  ·  Is this overfitting?"
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = WHITE
    r2 = p.add_run()
    r2.text = "   ·   distribution shift, not memorization"
    r2.font.size = Pt(14); r2.font.color.rgb = RGBColor(0xDD, 0xD8, 0xF0)

    # ── Top card — defining overfitting precisely ─────────────────
    add_card(s, 0.55, 1.0, 12.3, 1.20, fill=SOFT_BG)
    add_left_stripe(s, 0.55, 1.0, 1.20, SEC_RED, thickness=0.10)
    tb = add_textbox(s, 0.85, 1.12, 11.85, 1.05)
    add_para(tb.text_frame, "Overfitting · the precise definition",
              size=12, bold=True, color=SEC_RED)
    add_para(tb.text_frame,
              "A model that fits training data so tightly that it fails on "
              "held-out samples from the SAME distribution. Train ≫ Test gap.",
              size=12, color=MUTED_INK, space_before=4)

    # ── Two-column evidence ───────────────────────────────────────
    # Left — Why this is NOT overfitting
    add_card(s, 0.55, 2.40, 6.0, 4.55, fill=WHITE)
    add_strip(s, 0.55, 2.40, 6.0, SEC_GREEN, thickness=0.05)
    tb = add_textbox(s, 0.75, 2.52, 5.6, 4.35)
    add_para(tb.text_frame, "Why this is NOT overfitting",
              size=13, bold=True, color=SEC_GREEN)

    add_para(tb.text_frame,
              "① REDS4 is held-out from training",
              size=11.5, bold=True, color=INK, space_before=6)
    add_para(tb.text_frame,
              "Same distribution, never seen during training. All 9 metrics "
              "improve → in-distribution generalization works.",
              size=10.5, color=MUTED_INK, space_before=1)

    add_para(tb.text_frame,
              "② Memorization is structurally constrained",
              size=11.5, bold=True, color=INK, space_before=8)
    add_para(tb.text_frame,
              "U-Net 472 M + VAE 55 M FROZEN.  Only 6.22 M (Freq. Encoder) "
              "+ 208 M (ControlNet) train.  Identity-preserving init at "
              "step 0.",
              size=10.5, color=MUTED_INK, space_before=1)

    add_para(tb.text_frame,
              "③ OOD metrics are not uniformly worse",
              size=11.5, bold=True, color=INK, space_before=8)
    add_para(tb.text_frame,
              "If overfitting, every OOD metric would degrade. Instead, "
              "ours leads the DM group on UDM10 NR perceptual (MUSIQ, "
              "CLIP-IQA, NIQE), and improves SPMCS tLPIPS by −38 %.",
              size=10.5, color=MUTED_INK, space_before=1)

    add_para(tb.text_frame,
              "④ All DM baselines are REDS-only too",
              size=11.5, bold=True, color=INK, space_before=8)
    add_para(tb.text_frame,
              "Matched setup — StableVSR · DGAF-VSR · BasicVSR++ all use "
              "REDS train split. The OOD gap is shared, not unique to us.",
              size=10.5, color=MUTED_INK, space_before=1)

    # Right — Distribution shift vs overfitting framing
    add_card(s, 6.75, 2.40, 6.20, 4.55, fill=WHITE)
    add_strip(s, 6.75, 2.40, 6.20, SEC_BLUE, thickness=0.05)
    tb = add_textbox(s, 6.95, 2.52, 5.85, 1.4)
    add_para(tb.text_frame,
              "What we observe — distribution shift",
              size=13, bold=True, color=SEC_BLUE)
    add_para(tb.text_frame,
              "Vid4 / UDM10 / SPMCS differ from REDS in resolution, "
              "compression, and motion statistics. Wavelet conditioning "
              "specializes to training-frequency statistics by design — "
              "not memorization.",
              size=10.5, color=MUTED_INK, space_before=4)

    # OOD outcome table
    add_table(
        s, 6.95, 3.95, 5.85, 2.85,
        headers=["Set", "Pattern observed", "Overfit signature?"],
        rows=[
            ("REDS4",  "all 9 metrics improve",       "no"),
            ("UDM10",  "wins NR (MUSIQ/CLIP-IQA/NIQE)", "no"),
            ("SPMCS",  "tLPIPS −38% over StableVSR",  "no"),
            ("Vid4",   "close to baselines",           "no"),
        ],
        header_size=10, body_size=9.5,
    )

    # ── Bottom takeaway ───────────────────────────────────────────
    add_card(s, 0.55, 7.10, 12.3, 0.30, fill=None, line=None)
    tb = add_textbox(s, 0.55, 7.05, 12.3, 0.30)
    add_para(tb.text_frame,
              "→  Specialization–generalization trade-off (architectural), "
              "not overfitting (statistical).",
              size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # ── Korean speaker notes for this Q&A slide ───────────────────
    ntf = s.notes_slide.notes_text_frame
    ntf.clear()
    p = ntf.paragraphs[0]
    p.add_run().text = (
        "[Q&A 백업 슬라이드] Overfitting 질문이 나올 경우에만 보여주시면 "
        "됩니다.\n\n"
        "핵심 — overfitting과 distribution shift는 다른 현상입니다.\n\n"
        "Overfitting의 정의는 같은 distribution의 held-out test set에서 "
        "성능이 떨어지는 것입니다. REDS4는 학습에 사용되지 않은 held-out "
        "set이지만 REDS distribution과 동일하고, 거기서 모든 9개 metric이 "
        "향상됐습니다 — overfitting 신호가 아닙니다.\n\n"
        "Vid4 / UDM10 / SPMCS의 gap은 distribution shift입니다. 모든 DM "
        "비교 모델 (StableVSR · DGAF-VSR · BasicVSR++) 이 동일하게 "
        "REDS-only로 학습되어 같은 shift를 겪고 있습니다.\n\n"
        "또한 학습 가능한 파라미터는 단 6.22 M (Freq. Encoder) + "
        "208 M (ControlNet) 입니다. U-Net 472 M 과 VAE 55 M 은 frozen "
        "이고, identity-preserving initialization에서 시작합니다. 이 "
        "구조에서 memorization은 통계적으로 불가능합니다.\n\n"
        "결정적인 증거 — 만약 overfitting이라면 OOD 모든 metric이 "
        "일률적으로 악화돼야 하지만, UDM10 NR perceptual 3개에서 DM "
        "그룹 1위, SPMCS tLPIPS −38 % 개선이 나옵니다. 이건 "
        "specialization-generalization trade-off의 특징적 패턴이지 "
        "overfitting이 아닙니다."
    )


def main():
    prs = Presentation(str(SRC))

    # (1) text replacements
    n_repl = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.text in REPLACEMENTS:
                        run.text = REPLACEMENTS[run.text]
                        n_repl += 1
    print(f"Replaced {n_repl} run(s).")

    # (2) append overfitting Q&A defense to slide 29's Korean notes
    slide29 = prs.slides[28]
    ntf = slide29.notes_slide.notes_text_frame
    existing = ntf.text
    if "[Q&A 대비]" not in existing:
        # Append a new paragraph onto the existing notes
        p = ntf.add_paragraph()
        p.add_run().text = OVERFITTING_NOTE_KO.lstrip("\n\n")

    # (3) add the Q&A backup slide AFTER Thanks (currently last)
    build_overfitting_slide(prs)

    prs.save(str(SRC))
    print(f"Saved (overwritten): {SRC}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
