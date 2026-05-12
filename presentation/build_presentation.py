"""Build v3 of the Master's thesis defense deck.

Changes vs. v2 (author review round 3):
  1. Cover — thesis title centred horizontally.
  2. Speaker notes — rewritten in Korean; technical keywords stay English.
  3. Contents — section text now in black (subsections in muted gray).
  4. Removed the section-number 'page-number-like' badges at the top of
     every content slide (both the leading "1  " inside the title and the
     right-corner oval); the layout's auto slide number at bottom-right
     stays.
  5. Card shapes refined — softer fills, accent strips, less border weight.
  6. Added matplotlib-generated diagrams (band_pyramid, unet_injection,
     mag_phase_pipe) inside the Method section.
  7. Related Work rewritten with explicit motivation for *why* this method
     was designed, and includes recent 2025 papers (DiffVSR, STAR, DC-VSR,
     DLoRAL, UltraVSR, SeedVR2).
  8. Method · Why DT-CWT slide now follows the thesis's 3-reason structure
     (near shift-invariance · six directional subbands · explicit phase).
  9. Preliminaries section (Chapter 2) added — VSR & diffusion-based VSR
     formulation, DWT limitations, DT-CWT, Spatial Feature Transform.

Total slides: 31.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt, Inches


TEMPLATE = Path("/tmp/thesis_work/template.pptx")
FIG_DIR = Path("/tmp/thesis_work/fig_png")
OUT = Path("/home/user/StableVSR/presentation/WC_BD_SFT_Defense_Cho_YuJin.pptx")


# ── Palette ────────────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x12, 0x1E, 0x3D)        # near-black ink
MUTED_INK = RGBColor(0x3D, 0x3D, 0x4A)
GRAY = RGBColor(0x6B, 0x6B, 0x78)
LIGHT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
SOFT_BG = RGBColor(0xF6, 0xF4, 0xFB)
DEEP = RGBColor(0x15, 0x12, 0x7C)
SOFT_BG2 = RGBColor(0xF9, 0xF9, 0xFC)

# Section accent colors — distinct hues, cohesive saturation
SEC = {
    1: RGBColor(0x2C, 0x55, 0xA8),    # Introduction — blue
    2: RGBColor(0x0F, 0x86, 0x8E),    # Preliminaries — teal
    3: RGBColor(0x6A, 0x2E, 0xA6),    # Method — purple
    4: RGBColor(0xC5, 0x8A, 0x12),    # Experiments — gold
    5: RGBColor(0xCB, 0x42, 0x35),    # Results & Ablation — coral
    6: RGBColor(0x1F, 0x6E, 0x5C),    # Conclusion — forest
}
HIGH_BAND = SEC[1]
LOW_BAND = SEC[5]


# ── XML helper ──────────────────────────────────────────────────────
def _remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def add_content_slide(prs, section: int):
    """Add a slide from layout 12 and strip the default body placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == 10:
            _remove_shape(ph)
    return slide, SEC[section]


def set_section_title(slide, section: int, title: str, subtitle: str = ""):
    """White title text on the layout's dark navy bar.

    No section-number badges (per author feedback) — section identity is
    carried by the accent color of the cards on the slide.
    """
    ph = None
    for cand in slide.placeholders:
        if cand.placeholder_format.idx == 1:
            ph = cand
            break
    if ph is None:
        ph = slide.shapes.add_textbox(Inches(0.45), Inches(0.0),
                                       Inches(12.44), Inches(0.6))

    tf = ph.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0)
    p.space_after = Pt(0)

    main = p.add_run()
    main.text = title
    main.font.size = Pt(22)
    main.font.bold = True
    main.font.color.rgb = WHITE

    if subtitle:
        sub = p.add_run()
        sub.text = f"   ·   {subtitle}"
        sub.font.size = Pt(14)
        sub.font.bold = False
        sub.font.color.rgb = RGBColor(0xDD, 0xD8, 0xF0)


# ── Generic primitives ─────────────────────────────────────────────
def add_textbox(slide, left, top, width, height, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    return box


def add_para(tf, text, *, size=13, bold=False, italic=False, color=INK,
              align=PP_ALIGN.LEFT, space_before=2, space_after=4):
    if not tf.text and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    if color is not None:
        r.font.color.rgb = color
    return p


def add_card(slide, left, top, width, height, *, fill=WHITE, line=None,
              line_width=0.5, rounded=True, shadow=False):
    style = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        style, Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    shape.text_frame.text = ""
    return shape


def add_accent_strip(slide, left, top, width, color, *, thickness=0.045):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(thickness),
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_left_stripe(slide, left, top, height, color, *, thickness=0.08):
    """Vertical accent stripe on the left of a card."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(thickness), Inches(height),
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_picture_fit(slide, path, left, top, width, height, *, center=True):
    from PIL import Image
    with Image.open(path) as im:
        w_px, h_px = im.size
    ai = w_px / h_px
    ab = width / height
    if ai > ab:
        new_w = width
        new_h = width / ai
    else:
        new_h = height
        new_w = height * ai
    if center:
        left = left + (width - new_w) / 2
        top = top + (height - new_h) / 2
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                     Inches(new_w), Inches(new_h))


def add_table(slide, left, top, width, height, headers, rows, *,
              header_size=10.5, body_size=10, header_fill=DEEP,
              header_color=WHITE, zebra=True, highlight_row_idx=None,
              highlight_fill=RGBColor(0xFF, 0xF4, 0xCC)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left),
                                        Inches(top), Inches(width),
                                        Inches(height))
    tbl = tbl_shape.table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.margin_left = Inches(0.04)
        cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        r.font.size = Pt(header_size)
        r.font.bold = True
        r.font.color.rgb = header_color
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if highlight_row_idx is not None and i == highlight_row_idx:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_fill
            elif zebra and i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF7, 0xF7, 0xFB)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            tf.text = ""
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(body_size)
            r.font.color.rgb = INK
            if highlight_row_idx is not None and i == highlight_row_idx:
                r.font.bold = True
    return tbl_shape


def set_notes(slide, text: str):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    parts = [p.strip() for p in text.strip().split("\n\n") if p.strip()]
    for idx, part in enumerate(parts):
        p = notes_tf.paragraphs[0] if idx == 0 else notes_tf.add_paragraph()
        r = p.add_run()
        r.text = part
        r.font.size = Pt(11)


# ── Re-usable "stripe card" — modern card with thin colored stripe ──
def stripe_card(slide, left, top, width, height, *, accent, fill=WHITE,
                 line=LIGHT_GRAY, stripe_top=True, stripe_left=False):
    """White card with subtle border and a colored accent stripe."""
    add_card(slide, left, top, width, height, fill=fill,
              line=line, line_width=0.5)
    if stripe_top:
        add_accent_strip(slide, left, top, width, accent, thickness=0.05)
    if stripe_left:
        add_left_stripe(slide, left, top, height, accent, thickness=0.07)


def big_stat(slide, left, top, width, height, *, label, value, sub,
             accent):
    """KPI-style stat block."""
    stripe_card(slide, left, top, width, height, accent=accent)
    tb = add_textbox(slide, left + 0.1, top + 0.18, width - 0.2, height - 0.28)
    add_para(tb.text_frame, label, size=11, bold=True, color=accent,
              align=PP_ALIGN.CENTER, space_after=2)
    add_para(tb.text_frame, value, size=30, bold=True, color=INK,
              align=PP_ALIGN.CENTER, space_before=4, space_after=4)
    add_para(tb.text_frame, sub, size=10.5, color=GRAY,
              align=PP_ALIGN.CENTER, space_before=0)


# ──────────────────────────────────────────────────────────────────────────
# Cover & contents
# ──────────────────────────────────────────────────────────────────────────
def update_cover(prs):
    """Center the thesis title on the cover slide."""
    slide = prs.slides[0]
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        # Detect either v2 title or any leftover "SEMINAR"
        if "Wavelet-Conditioned" in txt or txt == "SEMINAR":
            tf = sh.text_frame
            tf.clear()
            tf.word_wrap = True
            # Centred, two-line title
            line1 = "Wavelet-Conditioned Band-Decoupled Spatial Feature Transform"
            line2 = "for Diffusion-Based Video Super-Resolution"
            p1 = tf.paragraphs[0]
            p1.alignment = PP_ALIGN.CENTER
            r1 = p1.add_run()
            r1.text = line1
            r1.font.size = Pt(26)
            r1.font.bold = True
            r1.font.color.rgb = INK
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(6)
            r2 = p2.add_run()
            r2.text = line2
            r2.font.size = Pt(26)
            r2.font.bold = True
            r2.font.color.rgb = INK
            # Recenter the textbox across the slide width.
            sh.left = Inches(0.5)
            sh.width = Inches(12.33)
            break

    set_notes(slide, (
        "안녕하십니까. IT공학과 조유진입니다.\n\n"
        "오늘 발표드릴 석사학위논문의 제목은 'Wavelet-Conditioned "
        "Band-Decoupled Spatial Feature Transform for Diffusion-Based "
        "Video Super-Resolution' 입니다.\n\n"
        "본 연구는 사전학습된 image diffusion model의 generative prior를 "
        "보존하면서도, 비디오 초해상화(VSR)에서 발생하는 temporal "
        "flickering 문제를 해결하기 위해, 주파수 도메인 조건화에 기반한 "
        "parameter-efficient한 어댑테이션 프레임워크를 제안합니다.\n\n"
        "전체 발표는 약 20분 정도 소요될 예정이며, 발표 후 질의응답 "
        "시간에 자유롭게 질문 부탁드립니다."
    ))


def update_contents(prs):
    """Rebuild Contents — section titles in BLACK, subsections in gray."""
    slide = prs.slides[1]
    target = None
    for sh in slide.shapes:
        if sh.name == "TextBox 5" and sh.has_text_frame:
            target = sh
            break
    if target is None:
        return

    target.left = Inches(5.41)
    target.top = Inches(2.40)
    target.width = Inches(7.40)
    target.height = Inches(5.0)

    tf = target.text_frame
    tf.clear()
    tf.word_wrap = True

    sections = [
        ("1.  Introduction",
         "Motivation · Related Work · Contributions"),
        ("2.  Preliminaries",
         "VSR · Diffusion · DWT → DT-CWT · SFT"),
        ("3.  Method",
         "Architecture · Why DT-CWT · SubbandBlock · BD-SFT · Loss"),
        ("4.  Experiments",
         "Datasets & Setup · Evaluation Metrics"),
        ("5.  Results & Ablation",
         "REDS4 · Cross-Dataset · Frequency · Band Injection · Discussion"),
        ("6.  Conclusion",
         "Limitations · Key Findings · Future Work"),
    ]

    for idx, (title, sub) in enumerate(sections):
        p_t = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p_t.alignment = PP_ALIGN.LEFT
        p_t.space_before = Pt(10 if idx > 0 else 0)
        p_t.space_after = Pt(2)
        rt = p_t.add_run()
        rt.text = title
        rt.font.size = Pt(18)
        rt.font.bold = True
        rt.font.color.rgb = INK  # ← black per author feedback

        p_s = tf.add_paragraph()
        p_s.alignment = PP_ALIGN.LEFT
        p_s.space_before = Pt(0)
        p_s.space_after = Pt(0)
        rs = p_s.add_run()
        rs.text = "         " + sub
        rs.font.size = Pt(11)
        rs.font.italic = True
        rs.font.color.rgb = GRAY

    set_notes(slide, (
        "발표는 총 여섯 개의 섹션으로 구성됩니다.\n\n"
        "1장 Introduction에서는 본 연구의 motivation, 관련 연구의 한계, "
        "그리고 본 논문의 contribution을 소개합니다.\n\n"
        "2장 Preliminaries에서는 VSR과 diffusion 기반 VSR의 기본 "
        "formulation, DWT의 한계와 DT-CWT의 동기, 그리고 SFT modulation의 "
        "기본 개념을 다룹니다.\n\n"
        "3장 Method에서는 제안하는 WC-BD-SFT 프레임워크 — 즉 Wavelet "
        "Conditioning Module, SubbandBlock, asymmetric BD-SFT injection, "
        "그리고 band-decoupled wavelet loss — 를 자세히 설명합니다.\n\n"
        "4-5장에서는 실험 셋업, REDS4 in-domain 평가, cross-dataset 평가, "
        "ablation, 그리고 trade-off에 대한 논의를 다루고,\n\n"
        "마지막 6장에서는 limitation과 future work으로 마무리합니다."
    ))


# ══════════════════════════════════════════════════════════════════════════
# Section 1 — Introduction
# ══════════════════════════════════════════════════════════════════════════
def slide_motivation(prs):
    s, accent = add_content_slide(prs, section=1)
    set_section_title(s, 1, "Motivation",
                       subtitle="Why VSR + Diffusion, and why is it hard?")

    # Top — three brief paradigm cards
    cards = [
        ("Demand for HR video",
         "고품질 HR 영상 수요 급증 · VSR이 spatial · temporal 정보로 HR 복원."),
        ("Diffusion priors win on detail",
         "CNN/GAN의 over-smoothing · instability를 넘어 realistic texture (SR3, DDPM)."),
        ("But: temporal flickering",
         "프레임 독립 stochastic denoising ⇒ HF detail이 frame 간 불일치 → flicker."),
    ]
    for i, (header, body) in enumerate(cards):
        x = 0.55 + i * 4.16
        stripe_card(s, x, 1.0, 4.0, 1.6, accent=accent)
        tb = add_textbox(s, x + 0.2, 1.15, 3.65, 1.45)
        add_para(tb.text_frame, header, size=13, bold=True, color=accent)
        add_para(tb.text_frame, body, size=11, color=MUTED_INK, space_before=4)

    # Middle — flickering visualization diagram
    add_picture_fit(s, FIG_DIR / "flickering_demo.png",
                     0.55, 2.85, 12.3, 2.6)

    # Bottom — research question
    add_card(s, 0.55, 5.7, 12.3, 1.3, fill=SOFT_BG2, line=LIGHT_GRAY,
              line_width=0.5)
    add_left_stripe(s, 0.55, 5.7, 1.3, accent, thickness=0.10)
    tb = add_textbox(s, 0.85, 5.83, 11.85, 1.15)
    add_para(tb.text_frame, "Research question",
              size=12, bold=True, color=accent)
    add_para(tb.text_frame,
              "사전학습된 image diffusion model을 parameter-efficient하게 video에 "
              "adapting하면서도, 3D attention · full fine-tuning 없이 temporal "
              "consistency를 달성할 수 있을까?",
              size=13, bold=True, color=INK, space_before=4)

    set_notes(s, (
        "첫 슬라이드 — motivation입니다.\n\n"
        "세 가지 흐름으로 정리하겠습니다. 첫째, 고해상도 영상 콘텐츠의 "
        "수요는 계속 증가하고 있습니다. VSR은 저해상도 프레임의 spatial "
        "정보와 인접 프레임의 temporal 정보를 함께 활용해 HR을 "
        "복원합니다.\n\n"
        "둘째, 최근 diffusion-based generative prior가 CNN이나 GAN 기반 "
        "방법보다 훨씬 사실적인 texture를 합성합니다. 자연스럽게 "
        "image-trained diffusion prior를 video로 가져오는 시도가 "
        "활발해졌습니다.\n\n"
        "셋째, 그러나 결정적인 문제가 있습니다 — 프레임마다 독립적으로 "
        "stochastic denoising을 수행하기 때문에, 복원되는 high-frequency "
        "detail이 프레임 간에 일관되지 않습니다. 결과적으로 영상에서 "
        "texture가 어른거리는 temporal flickering 현상이 발생합니다.\n\n"
        "본 연구의 research question은 이렇습니다 — 가벼운 모듈만으로, "
        "사전학습된 image diffusion model을 video에 temporally consistent "
        "하게 적응시킬 수 있을까? 제 답은 'pixel을 warping하지 말고, "
        "shift-stable한 frequency-domain prior로 conditioning하자'는 "
        "것입니다."
    ))


def slide_related_work_paradigms(prs):
    s, accent = add_content_slide(prs, section=1)
    set_section_title(s, 1, "Related Work · VSR Paradigms",
                       subtitle="From CNN regression to diffusion priors")

    # Top — timeline diagram
    add_picture_fit(s, FIG_DIR / "vsr_timeline.png",
                     0.40, 0.95, 12.55, 3.0)

    # Bottom — three paradigm summary cards
    CNN_COL = RGBColor(0x7C, 0x7C, 0x8E)
    GAN_COL = RGBColor(0xD8, 0x71, 0x58)
    cards = [
        ("CNN-based",
         ["Alignment-driven: TOFlow / EDVR / BasicVSR++ / VRT.",
          "Pixel loss ⇒ over-smoothed · alignment fragile."],
         CNN_COL),
        ("GAN-based",
         ["SRGAN · ESRGAN · VideoGigaGAN (CVPR'24).",
          "Training instability · bounded texture diversity."],
         GAN_COL),
        ("Diffusion-based",
         ["Image-prior adapt: StableVSR (baseline) · DGAF-VSR.",
          "Video-native: Upscale-A-Video · SVD.",
          "Recent 2025: DiffVSR / STAR / DC-VSR / DLoRAL / UltraVSR / SeedVR2."],
         HIGH_BAND),
    ]
    for i, (h, lines, col) in enumerate(cards):
        x = 0.55 + i * 4.16
        stripe_card(s, x, 4.15, 4.0, 2.85, accent=col)
        tb = add_textbox(s, x + 0.2, 4.30, 3.65, 2.65)
        add_para(tb.text_frame, h, size=13, bold=True, color=col)
        for line in lines:
            add_para(tb.text_frame, "▸ " + line, size=10.5,
                      color=MUTED_INK, space_before=3)

    set_notes(s, (
        "관련 연구를 두 축으로 정리했습니다.\n\n"
        "왼쪽은 CNN과 GAN 기반 VSR 입니다. CNN 계열은 alignment에서 "
        "출발해 — TOFlow의 explicit optical flow부터, EDVR의 deformable "
        "convolution, BasicVSR과 BasicVSR++의 bidirectional propagation, "
        "그리고 VRT와 RVRT 같은 transformer 기반 방법까지 발전했습니다. "
        "하지만 pixel-wise loss를 사용하기 때문에 결과가 over-smoothing "
        "되는 한계가 있습니다. GAN은 더 sharp한 texture를 만들지만 "
        "training instability와 제한된 texture diversity 문제가 "
        "있습니다.\n\n"
        "오른쪽은 diffusion 기반입니다. 두 방향이 있는데, 첫 번째는 "
        "사전학습된 image diffusion prior에 가벼운 temporal module을 "
        "더하는 방식 — StableVSR가 대표적이고 이 논문의 직접 baseline 입니다. "
        "ControlNet에 RAFT로 warping한 이웃 프레임의 x̂_0을 넣어주는 "
        "구조죠. DGAF-VSR는 가장 최근(CVPR 2026) 비교 대상으로, "
        "feature-domain warping을 통해 temporal guidance를 줍니다.\n\n"
        "두 번째는 video-native diffusion을 처음부터 학습하는 방식 — "
        "Stable Video Diffusion이 대표적인데 수십억 파라미터의 비용이 "
        "듭니다. 2025년에는 robustness 방향(DiffVSR, STAR, DC-VSR)과 "
        "efficiency 방향(DLoRAL, UltraVSR, SeedVR2)의 연구가 활발히 "
        "진행되었습니다."
    ))


def slide_related_work_gap(prs):
    """The 'why' slide — frequency-domain gap → our motivation."""
    s, accent = add_content_slide(prs, section=1)
    set_section_title(s, 1, "Related Work · Where the Gap Is",
                       subtitle="Why frequency-domain conditioning?")

    # Top: frequency-domain prior works
    stripe_card(s, 0.55, 1.0, 12.3, 2.0, accent=accent)
    tb = add_textbox(s, 0.85, 1.15, 11.85, 1.85)
    add_para(tb.text_frame,
              "Frequency-domain methods exist — but with two key limitations",
              size=14, bold=True, color=accent)
    add_para(tb.text_frame,
              "▸  Image-domain : MWCNN · WaveletSRNet · DWSR — frequency loss "
              "improves SR sharpness.",
              size=12, color=MUTED_INK, space_before=4)
    add_para(tb.text_frame,
              "▸  Video      : FTVSR — frequency-aware temporal attention for "
              "compressed VSR.",
              size=12, color=MUTED_INK, space_before=2)
    add_para(tb.text_frame,
              "▸  Diffusion-SR : ResQu (quaternion wavelets) · WaveDiT "
              "(wavelet-spectrum DiT) — single image only.",
              size=12, color=MUTED_INK, space_before=2)

    # Middle two cards — the limitations
    stripe_card(s, 0.55, 3.15, 6.0, 2.0, accent=LOW_BAND, stripe_top=True)
    tb = add_textbox(s, 0.75, 3.30, 5.7, 1.85)
    add_para(tb.text_frame, "① Limited to single-image SR",
              size=12, bold=True, color=LOW_BAND)
    add_para(tb.text_frame,
              "Video temporal consistency 문제를 다루지 않음. "
              "Diffusion 기반 frequency 연구들은 모두 image SR만 대상.",
              size=11, color=MUTED_INK, space_before=4)

    stripe_card(s, 6.85, 3.15, 6.0, 2.0, accent=LOW_BAND, stripe_top=True)
    tb = add_textbox(s, 7.05, 3.30, 5.7, 1.85)
    add_para(tb.text_frame, "② Use shift-variant DWT",
              size=12, bold=True, color=LOW_BAND)
    add_para(tb.text_frame,
              "Critical decimation 때문에 sub-pixel motion에도 wavelet "
              "coefficient가 크게 흔들림 → video flickering 악화.",
              size=11, color=MUTED_INK, space_before=4)

    # Bottom — our proposed angle, prominent
    add_card(s, 0.55, 5.35, 12.3, 1.6, fill=SOFT_BG2, line=accent,
              line_width=1.0)
    add_left_stripe(s, 0.55, 5.35, 1.6, accent, thickness=0.10)
    tb = add_textbox(s, 0.85, 5.45, 11.85, 1.45)
    add_para(tb.text_frame, "→  Our motivation",
              size=12, bold=True, color=accent)
    add_para(tb.text_frame,
              "Image SR이 아닌 video VSR을 대상으로, DWT가 아닌 "
              "near-shift-invariant DT-CWT를 사용해, ",
              size=13, bold=True, color=INK, space_before=4)
    add_para(tb.text_frame,
              "사전학습된 diffusion prior를 freeze한 채 asymmetric "
              "encoder-decoder injection으로 conditioning한다.",
              size=13, bold=True, color=INK)

    set_notes(s, (
        "이 슬라이드가 가장 중요합니다 — 왜 제가 이 알고리즘을 "
        "설계했는지에 대한 답입니다.\n\n"
        "Frequency-domain을 활용한 연구는 이미 많이 있습니다. "
        "Image-domain에서는 MWCNN, WaveletSRNet, DWSR 같은 wavelet "
        "기반 손실이 super-resolution의 sharpness를 높여줍니다. Video "
        "쪽에서는 FTVSR가 compressed video를 위한 frequency-aware "
        "attention을 제안했고, 가장 최근에는 diffusion-SR 분야에서 "
        "ResQu와 WaveDiT가 wavelet conditioning을 적용했습니다.\n\n"
        "그런데 두 가지 핵심적인 한계가 있습니다.\n\n"
        "첫째, diffusion 기반 frequency 연구들은 모두 single-image SR을 "
        "대상으로 합니다. Video VSR의 temporal consistency 문제는 "
        "다루지 않습니다.\n\n"
        "둘째 — 그리고 이게 결정적인 부분인데 — 기존 방법들은 모두 "
        "standard DWT 또는 그 변형을 씁니다. 그런데 DWT는 critical "
        "decimation 때문에 shift-variant 합니다. 인접 프레임 사이의 "
        "sub-pixel motion에도 wavelet coefficient가 크게 흔들리죠. "
        "이건 오히려 video flickering을 악화시킬 수 있습니다.\n\n"
        "그래서 제 motivation은 명확합니다 — 첫째, single-image SR이 "
        "아닌 video VSR을 대상으로 하고, 둘째, DWT가 아닌 near-"
        "shift-invariant 한 DT-CWT를 사용하고, 셋째, 사전학습된 "
        "diffusion prior를 freeze한 채 가볍게 conditioning하기 위해 "
        "asymmetric encoder-decoder injection을 설계합니다."
    ))


def slide_contributions(prs):
    s, accent = add_content_slide(prs, section=1)
    set_section_title(s, 1, "Contributions")

    contribs = [
        ("WC-BD-SFT framework",
         ["DT-CWT 기반 frequency-domain conditioning을 frozen pre-trained "
          "diffusion U-Net에 주입.",
          "Asymmetric encoder-decoder injection — HIGH → decoder (texture), "
          "LOW → encoder (structure)."]),
        ("Decoupled magnitude-phase processing",
         ["Trigonometric phase encoding — 2π wraparound 불연속 회피.",
          "Frequency Encoder 6.22 M (≈3% of ControlNet) · "
          "identity-preserving init."]),
        ("Band-decoupled wavelet loss",
         ["HIGH는 magnitude only, LOW는 magnitude-weighted angular distance.",
          "Latent ε-MSE에 더해 pixel-space frequency supervision."]),
        ("Strong empirical gains (REDS4)",
         ["−62.9 % tLPIPS (41.11 → 15.25) · −37.5 % LPIPS",
          "+103 % high-frequency spectral power retention vs. StableVSR."]),
    ]
    for i, (header, lines) in enumerate(contribs):
        col = i % 2
        row = i // 2
        x = 0.55 + col * 6.15
        y = 1.0 + row * 2.95
        stripe_card(s, x, y, 6.0, 2.75, accent=accent)
        tb = add_textbox(s, x + 0.2, y + 0.15, 5.65, 2.55)
        add_para(tb.text_frame, header, size=14, bold=True, color=accent)
        for line in lines:
            add_para(tb.text_frame, "▸  " + line, size=11.5,
                      color=MUTED_INK, space_before=4)

    set_notes(s, (
        "본 연구의 contribution은 네 가지입니다.\n\n"
        "첫째 — WC-BD-SFT 프레임워크 자체입니다. 핵심은 asymmetric "
        "band-decoupled injection — 고주파 wavelet band는 decoder로 보내 "
        "texture 합성에 기여하고, 저주파 band는 encoder로 보내 structural "
        "feature 추출에 기여하도록 분리한 점입니다.\n\n"
        "둘째 — magnitude와 phase를 분리해 처리합니다. Phase는 sin과 "
        "cos pair로 인코딩해 ±π 경계에서의 wraparound 불연속을 "
        "회피합니다. 전체 Frequency Encoder는 6.22 M 파라미터로 "
        "ControlNet의 약 3% 수준이며, identity-preserving "
        "initialization으로 학습 시작점에서 pre-trained prior를 "
        "그대로 보존합니다.\n\n"
        "셋째 — band-decoupled wavelet loss를 새롭게 도입했습니다. "
        "HIGH band는 magnitude만 supervise하여 자유로운 texture 합성을 "
        "허용하고, LOW band는 magnitude-weighted angular distance로 "
        "구조의 phase 정확도까지 강제합니다.\n\n"
        "넷째 — 실험 결과입니다. Matched REDS-only 학습 조건에서, "
        "tLPIPS는 63% 감소, LPIPS는 37.5% 감소, 그리고 high-frequency "
        "spectral power retention은 StableVSR 대비 두 배 이상 향상됩니다."
    ))


# ══════════════════════════════════════════════════════════════════════════
# Section 2 — Preliminaries (NEW)
# ══════════════════════════════════════════════════════════════════════════
def slide_prelim_vsr_diffusion(prs):
    s, accent = add_content_slide(prs, section=2)
    set_section_title(s, 2, "Preliminaries · VSR & Diffusion",
                       subtitle="Formulation and the StableVSR pipeline")

    # Top equation strip — VSR formulation
    add_card(s, 0.55, 0.95, 12.3, 0.85, fill=SOFT_BG2, line=LIGHT_GRAY,
              line_width=0.5)
    add_left_stripe(s, 0.55, 0.95, 0.85, accent, thickness=0.08)
    tb = add_textbox(s, 0.55, 1.02, 12.3, 0.78, anchor=MSO_ANCHOR.MIDDLE)
    add_para(tb.text_frame, "VSR — LR → HR with BIx4 degradation",
              size=11, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_para(tb.text_frame,
              "I^LR  =  𝒟_{bicubic}(I^HR ;  s = 4)",
              size=15, bold=True, color=INK, align=PP_ALIGN.CENTER,
              space_before=2)

    # Middle — diffusion process diagram
    add_picture_fit(s, FIG_DIR / "diffusion_process.png",
                     0.55, 1.95, 12.3, 2.85)

    # Bottom — StableVSR baseline card with key text
    stripe_card(s, 0.55, 5.00, 6.0, 1.95, accent=accent)
    L = add_textbox(s, 0.75, 5.15, 5.6, 1.80)
    add_para(L.text_frame, "StableVSR baseline (Rota et al., ECCV'24)",
              size=13, bold=True, color=accent)
    add_para(L.text_frame,
              "▸  ControlNet + RAFT-warped neighbor x̂₀",
              size=11, color=MUTED_INK, space_before=4)
    add_para(L.text_frame,
              "▸  Frame-wise Bidirectional Sampling",
              size=11, color=MUTED_INK, space_before=2)
    add_para(L.text_frame,
              "→  본 연구는 이 scaffold를 그대로 두고 WCM만 추가.",
              size=11, italic=True, color=accent, space_before=4)

    # Bottom-right — perception-distortion trade-off note
    stripe_card(s, 6.85, 5.00, 6.0, 1.95, accent=accent)
    R = add_textbox(s, 7.05, 5.15, 5.6, 1.80)
    add_para(R.text_frame, "Perception–distortion trade-off",
              size=13, bold=True, color=accent)
    add_para(R.text_frame,
              "Blau & Michaeli (2018) — pixel fidelity와 perceptual "
              "quality는 동시 최대화 불가.",
              size=11, color=MUTED_INK, space_before=4)
    add_para(R.text_frame,
              "본 연구는 perceptual 극단을 추구하는 설계.",
              size=11, italic=True, color=accent, space_before=4)

    set_notes(s, (
        "Preliminaries 첫 번째 슬라이드입니다. VSR과 diffusion-based "
        "VSR의 기본 개념을 정리합니다.\n\n"
        "왼쪽 — VSR은 LR 시퀀스에서 HR 시퀀스를 복원하는 task입니다. "
        "보통 bicubic ×4 downsampling으로 LR을 만들고, 두 핵심 "
        "컴포넌트가 temporal alignment와 temporal aggregation 입니다. "
        "한 가지 중요한 이론적 배경 — Blau와 Michaeli의 perception-"
        "distortion trade-off가 있습니다. Pixel fidelity와 perceptual "
        "quality를 동시에 최대화할 수 없다는 것이죠. 본 연구는 의도적으로 "
        "perceptual quality 극단을 추구합니다.\n\n"
        "오른쪽 — latent diffusion model 입니다. Forward 과정에서 latent에 "
        "noise를 점진적으로 추가하고, U-Net이 이 noise를 예측하도록 "
        "epsilon-MSE로 학습합니다. 본 연구는 Stable Diffusion v2.1을 "
        "backbone으로 사용합니다.\n\n"
        "StableVSR는 본 연구의 직접 baseline 입니다. ControlNet으로 SD를 "
        "VSR에 적응시키는 구조이고, Temporal Conditioning Module(TCM)이 "
        "핵심 입니다. 매 denoising step마다 인접 프레임의 x_hat_0 예측을 "
        "RAFT optical flow로 현재 프레임에 warping해 ControlNet input으로 "
        "사용합니다. 또한 Frame-wise Bidirectional Sampling으로 forward "
        "방향과 backward 방향을 번갈아 적용해 temporal consistency를 "
        "강화합니다. 본 연구는 이 scaffold를 변경 없이 그대로 사용하고, "
        "Wavelet Conditioning Module만 추가합니다."
    ))


def slide_prelim_dwt_dtcwt(prs):
    s, accent = add_content_slide(prs, section=2)
    set_section_title(s, 2, "Preliminaries · DWT → DT-CWT",
                       subtitle="From shift-variant to shift-invariant")

    # Left — DWT limitations
    stripe_card(s, 0.55, 1.0, 6.0, 3.6, accent=LOW_BAND, stripe_top=True)
    L = add_textbox(s, 0.75, 1.15, 5.55, 3.4)
    add_para(L.text_frame, "Discrete Wavelet Transform (DWT)",
              size=14, bold=True, color=LOW_BAND)
    add_para(L.text_frame,
              "✗  Shift-variance — critical decimation 때문에 sub-pixel "
              "shift에도 coefficient가 크게 흔들림.",
              size=12, color=MUTED_INK, space_before=6)
    add_para(L.text_frame,
              "✗  3 directional subbands per scale — horizontal · "
              "vertical · diagonal만.",
              size=12, color=MUTED_INK, space_before=2)
    add_para(L.text_frame,
              "✗  Real-valued — phase 정보 미포함.",
              size=12, color=MUTED_INK, space_before=2)

    # Right — DT-CWT
    stripe_card(s, 6.85, 1.0, 6.0, 3.6, accent=HIGH_BAND, stripe_top=True)
    R = add_textbox(s, 7.05, 1.15, 5.55, 3.4)
    add_para(R.text_frame, "Dual-Tree Complex Wavelet (DT-CWT)",
              size=14, bold=True, color=HIGH_BAND)
    add_para(R.text_frame,
              "✓  Near shift-invariance — 두 parallel tree (Hilbert transform pair).",
              size=11.5, color=MUTED_INK, space_before=4)
    add_para(R.text_frame,
              "✓  6 directional subbands per scale — ±15°, ±45°, ±75°.",
              size=11.5, color=MUTED_INK, space_before=2)
    add_para(R.text_frame,
              "✓  Complex coefficients C = a + ib  ⇒  (magnitude, phase) 분해:",
              size=11.5, color=MUTED_INK, space_before=2)
    add_para(R.text_frame,
              "M = |C| = √(a² + b²)     ·     φ = ∠C ∈ [−π, π)",
              size=12.5, bold=True, color=INK, space_before=4,
              align=PP_ALIGN.CENTER)

    # Bottom — empirical proof image
    add_picture_fit(s, FIG_DIR / "dtcwt_shift_zoom_v4-1.png",
                     0.55, 4.85, 12.3, 2.0)
    cap = add_textbox(s, 0.55, 6.85, 12.3, 0.20)
    add_para(cap.text_frame,
              "Empirically — DT-CWT magnitudes 5.09× more stable than DWT "
              "under sub-pixel shifts (σ = 0.019 vs. 0.087).",
              size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    set_notes(s, (
        "Preliminaries 두 번째 — DWT의 한계와 DT-CWT의 동기 입니다.\n\n"
        "왼쪽 — 표준 DWT의 세 가지 한계 입니다. 첫째, critical decimation "
        "때문에 shift-variant 합니다. 입력이 sub-pixel만큼만 움직여도 "
        "wavelet coefficient가 크게 흔들립니다. Video VSR 입장에서는 "
        "치명적인 단점입니다. 둘째, scale 당 directional subband가 "
        "horizontal, vertical, diagonal 셋 뿐입니다. Texture의 "
        "다양한 방향성을 잘 표현하지 못합니다. 셋째, real-valued라서 "
        "phase 정보가 없습니다.\n\n"
        "오른쪽 — DT-CWT는 이 세 가지를 모두 해결합니다. 첫째, Hilbert "
        "transform pair를 이루는 두 개의 parallel DWT tree로 구성되어 "
        "있어 magnitude가 near shift-invariant 합니다. 둘째, scale 당 "
        "방향 subband가 ±15°, ±45°, ±75° 의 여섯 방향으로 두 배입니다. "
        "셋째, complex-valued 라서 magnitude와 phase로 분해할 수 "
        "있습니다.\n\n"
        "아래 그림은 sub-pixel shift 실험 결과 입니다. DT-CWT magnitude의 "
        "per-pixel 표준편차가 DWT보다 5.09배 더 안정적입니다. 이 "
        "shift-stability가 본 연구에서 DT-CWT를 선택한 핵심 이유 입니다."
    ))


def slide_prelim_sft(prs):
    s, accent = add_content_slide(prs, section=2)
    set_section_title(s, 2, "Preliminaries · Spatial Feature Transform",
                       subtitle="Spatially-varying affine modulation")

    # Equation block — clean centered
    add_card(s, 2.5, 1.0, 8.4, 1.3, fill=SOFT_BG2, line=LIGHT_GRAY,
              line_width=0.5)
    add_left_stripe(s, 2.5, 1.0, 1.3, accent, thickness=0.08)
    tb = add_textbox(s, 2.5, 1.10, 8.4, 1.1, anchor=MSO_ANCHOR.MIDDLE)
    add_para(tb.text_frame, "Spatial Feature Transform (Wang et al., CVPR 2018)",
              size=11, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_para(tb.text_frame, "X′  =  X ⊙ γ  +  β",
              size=26, bold=True, color=INK, align=PP_ALIGN.CENTER,
              space_before=4)
    add_para(tb.text_frame, "feature를 channel-wise 가 아닌 spatial-wise로 "
              "affine modulation.",
              size=10.5, italic=True, color=GRAY, align=PP_ALIGN.CENTER,
              space_before=4)

    # Comparison
    cards = [
        ("AdaIN", "Globally uniform.\nSpatial detail 손실.", GRAY),
        ("Cross-attention", "Spatially varying.\nO(N²) cost.", GRAY),
        ("SFT (ours)", "Spatially varying + linear-time.\n"
                       "Wavelet conditioning에 최적.", accent),
    ]
    for i, (h, body, col) in enumerate(cards):
        x = 0.55 + i * 4.16
        stripe_card(s, x, 2.65, 4.0, 1.7, accent=col)
        tb = add_textbox(s, x + 0.18, 2.78, 3.65, 1.55)
        add_para(tb.text_frame, h, size=13, bold=True, color=col)
        add_para(tb.text_frame, body, size=11, color=MUTED_INK, space_before=4)

    # Why SFT for our task
    add_card(s, 0.55, 4.6, 12.3, 2.35, fill=SOFT_BG2, line=LIGHT_GRAY,
              line_width=0.5)
    add_left_stripe(s, 0.55, 4.6, 2.35, accent, thickness=0.10)
    tb = add_textbox(s, 0.85, 4.78, 11.85, 2.15)
    add_para(tb.text_frame, "Why SFT fits this work",
              size=13, bold=True, color=accent)
    add_para(tb.text_frame,
              "▸  Wavelet coefficient는 본질적으로 spatial하게 분포 — "
              "spatial-wise modulation이 자연스러운 매칭.",
              size=12, color=MUTED_INK, space_before=4)
    add_para(tb.text_frame,
              "▸  Cross-attention 대비 가벼움 — H × W 크기의 latent에서 "
              "quadratic cost 회피.",
              size=12, color=MUTED_INK, space_before=2)
    add_para(tb.text_frame,
              "▸  Identity-preserving init (γ=1, β=0)이 자연스럽게 정의됨 — "
              "학습 시작 시 baseline 동작 보존.",
              size=12, color=MUTED_INK, space_before=2)

    set_notes(s, (
        "Preliminaries 마지막 — Spatial Feature Transform 입니다.\n\n"
        "SFT는 Wang et al.이 CVPR 2018에서 image SR을 위해 제안한 "
        "modulation 기법 입니다. Feature map X에 대해 channel-wise가 "
        "아닌 spatial-wise로 affine 변환을 적용합니다 — γ로 곱하고 "
        "β를 더하는 것이죠.\n\n"
        "다른 modulation 기법과 비교하면, AdaIN은 globally uniform이라 "
        "spatial detail을 잃고, cross-attention은 spatially varying이지만 "
        "spatial size에 대해 quadratic cost가 듭니다. SFT는 spatially "
        "varying이면서 linear time입니다.\n\n"
        "SFT가 본 연구에 잘 맞는 이유는 세 가지 입니다. 첫째, wavelet "
        "coefficient 자체가 본질적으로 spatial한 분포를 갖기 때문에 "
        "spatial-wise modulation이 자연스러운 매칭입니다. 둘째, "
        "cross-attention보다 훨씬 가벼워서 latent의 H×W가 커도 "
        "감당 가능합니다. 셋째, identity-preserving initialization을 "
        "자연스럽게 정의할 수 있어서 — γ를 1, β를 0으로 두면 — 학습 "
        "시작 시점에 사전학습된 baseline이 그대로 보존됩니다."
    ))


# ══════════════════════════════════════════════════════════════════════════
# Section 3 — Method
# ══════════════════════════════════════════════════════════════════════════
def slide_why_dtcwt(prs):
    s, accent = add_content_slide(prs, section=3)
    set_section_title(s, 3, "Method · Why DT-CWT, not DWT or Fourier",
                       subtitle="Three reasons from the thesis")

    # Top — DT-CWT decomposition equation
    add_card(s, 0.55, 0.95, 12.3, 0.85, fill=SOFT_BG2, line=LIGHT_GRAY,
              line_width=0.5)
    add_left_stripe(s, 0.55, 0.95, 0.85, accent, thickness=0.08)
    tb = add_textbox(s, 0.55, 1.02, 12.3, 0.78, anchor=MSO_ANCHOR.MIDDLE)
    add_para(tb.text_frame, "DT-CWT decomposition (J = 4)",
              size=10.5, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_para(tb.text_frame,
              "DT-CWT_{J=4}(I^{LR}_t)  =  ( I^{LP}_t ,  { C^{(j,d)}_t }_{j=1..4, d=1..6} ),     "
              "C^{(j,d)} ∈ ℂ^{H/2ʲ × W/2ʲ × 3}",
              size=13, bold=True, color=INK, align=PP_ALIGN.CENTER,
              space_before=2)

    reasons = [
        ("① Near shift-invariance",
         "Magnitudes ≈ stable under sub-pixel translations.\n"
         "DWT는 critical decimation으로 frame 간 inconsistent conditioning.",
         HIGH_BAND),
        ("② Six directional subbands",
         "±15°, ±45°, ±75° at each scale (2× DWT's 3 subbands).\n"
         "Hair / foliage / fabric의 임의 방향 edge 보존 — oriented HF detail 합성에 적합.",
         accent),
        ("③ Complex coeffs → magnitude · phase",
         "M = |C| = √(a²+b²) → texture (shift-stable energy)\n"
         "φ = ∠C → structure (precise spatial localization) — BD-SFT pathway에 직접 대응.",
         LOW_BAND),
    ]
    for i, (h, body, col) in enumerate(reasons):
        y = 1.95 + i * 1.55
        stripe_card(s, 0.55, y, 12.3, 1.40, accent=col, stripe_left=True,
                     stripe_top=False)
        tb = add_textbox(s, 0.85, y + 0.10, 12.0, 1.25)
        add_para(tb.text_frame, h, size=13.5, bold=True, color=col)
        for line in body.split("\n"):
            add_para(tb.text_frame, line, size=11, color=MUTED_INK,
                      space_before=2, space_after=2)

    # Fourier note
    tb = add_textbox(s, 0.55, 6.65, 12.3, 0.35)
    add_para(tb.text_frame,
              "vs. Fourier — globally-supported basis ⇒ cannot vary spatially. "
              "DT-CWT combines Fourier's mag-phase factorization with wavelet's "
              "spatial localization.",
              size=10.5, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    set_notes(s, (
        "Method 섹션의 첫 슬라이드 — 왜 DT-CWT를 선택했는지에 대한 "
        "논문의 3가지 핵심 이유를 그대로 따라 설명합니다.\n\n"
        "첫째, near shift-invariance 입니다. DT-CWT magnitude는 sub-pixel "
        "translation에 대해 approximately stable합니다. 반면 DWT는 "
        "critical decimation 때문에 같은 shift에서도 coefficient가 크게 "
        "변합니다. Video VSR에서 인접 프레임은 sub-pixel motion으로 "
        "차이가 나므로, shift-variant한 표현은 frame 간 inconsistent "
        "conditioning을 만듭니다.\n\n"
        "둘째, six directional subbands per scale 입니다. DT-CWT는 각 "
        "scale에서 ±15°, ±45°, ±75° 의 6개 방향 subband를 제공합니다. "
        "DWT의 3개에 비해 두 배 풍부하죠. 머리카락, 잎사귀, 직물 같은 "
        "자연 texture의 임의 방향 edge를 잘 보존합니다. Decoder가 "
        "directionally coherent한 high-frequency detail을 합성하는 데 "
        "이 정보가 필수적입니다.\n\n"
        "셋째, complex coefficient의 magnitude-phase 분해 입니다. "
        "Magnitude는 shift-stable한 energy descriptor — texture 정보. "
        "Phase는 spatial localization — structure 정보. 이 분해가 본 "
        "연구의 band-decoupled SFT pathway에 직접 대응됩니다.\n\n"
        "Fourier transform도 complex coefficient를 주지만 globally-"
        "supported basis라서 spatial하게 conditioning을 변화시킬 수 "
        "없습니다. DT-CWT는 Fourier의 mag-phase factorization과 "
        "wavelet의 spatial localization을 동시에 갖춥니다."
    ))


def slide_band_partition(prs):
    s, accent = add_content_slide(prs, section=3)
    set_section_title(s, 3, "Method · Multi-Scale Decomposition",
                       subtitle="DT-CWT → band partitioning at j=2 / j=3")

    # The matplotlib diagram is the centerpiece
    add_picture_fit(s, FIG_DIR / "band_pyramid.png",
                     0.55, 1.0, 12.3, 3.4)

    # Rationale strip
    add_card(s, 0.55, 4.65, 12.3, 2.3, fill=SOFT_BG2, line=LIGHT_GRAY,
              line_width=0.5)
    add_left_stripe(s, 0.55, 4.65, 2.3, accent, thickness=0.10)
    tb = add_textbox(s, 0.85, 4.80, 11.85, 2.10)
    add_para(tb.text_frame, "Why this 2 : 2 split at j=2 / j=3?",
              size=13, bold=True, color=accent)
    add_para(tb.text_frame,
              "▸  Spectral — LR 신호 spectrum을 정확히 절반으로 분리 "
              "(texture 상위, structure 하위). 각 SFT pathway가 "
              "non-overlapping range에 특화.",
              size=11.5, color=MUTED_INK, space_before=4)
    add_para(tb.text_frame,
              "▸  Architectural — 2 scales × 256 ch = 512 ch ⇒ U-Net의 "
              "block_out_channels[1] = 512와 정확히 일치, 추가 projection 불필요.",
              size=11.5, color=MUTED_INK, space_before=2)
    add_para(tb.text_frame,
              "▸  J = 4 — j=4 subband의 H/16 × W/16 size가 SD-VAE 8× 와 "
              "U-Net 2× 다운샘플링 위에서 자연스럽게 정렬.",
              size=11.5, color=MUTED_INK, space_before=2)

    set_notes(s, (
        "본 연구의 multi-scale decomposition과 band partitioning 입니다.\n\n"
        "위쪽 그림 — LR 입력 프레임을 4-level DT-CWT로 분해하면 4개의 "
        "scale이 나옵니다. 각 scale은 6개의 directional subband를 "
        "갖습니다. 이 4개 scale을 j=2와 j=3 사이에서 두 그룹으로 "
        "나눕니다.\n\n"
        "HIGH band는 j=1과 j=2 — normalized frequency 1/8 ~ 1/2 cyc/px, "
        "fine texture와 sharp edge에 해당합니다. 이건 up_blocks[1] 즉 "
        "decoder로 들어갑니다. LOW band는 j=3과 j=4 — 0 ~ 1/8 cyc/px, "
        "coarse structure에 해당합니다. 이건 down_blocks[1] 즉 encoder로 "
        "들어갑니다. 별도로 real-valued lowpass component(LP)는 wavelet "
        "loss에서 reference로 사용됩니다.\n\n"
        "왜 이렇게 2대 2로 나눴는지 세 가지 이유가 있습니다.\n\n"
        "첫째, spectral 관점 — LR 신호의 주파수 spectrum을 정확히 "
        "절반으로 나눠 texture 상위와 structure 하위로 분리할 수 "
        "있습니다. 각 SFT pathway가 non-overlapping range에 특화됩니다.\n\n"
        "둘째, architectural 관점 — 2개 scale의 SFT output을 channel-wise "
        "concat하면 2 × 256 = 512 채널이 됩니다. 이건 Stable Diffusion "
        "U-Net의 block_out_channels[1]과 정확히 일치해서, 추가적인 "
        "channel projection 없이 그대로 modulation 가능합니다.\n\n"
        "셋째, J=4를 선택한 이유 — j=4 subband의 spatial size H/16 × "
        "W/16이 SD-VAE의 8× 다운샘플링과 U-Net의 2× 다운샘플링과 "
        "자연스럽게 정렬됩니다."
    ))


def slide_subbandblock_detail(prs):
    s, accent = add_content_slide(prs, section=3)
    set_section_title(s, 3, "Method · SubbandBlock",
                       subtitle="Magnitude / phase decoupled pipeline")

    add_picture_fit(s, FIG_DIR / "mag_phase_pipe.png",
                     0.55, 0.95, 12.3, 3.3)

    # Two callouts below
    stripe_card(s, 0.55, 4.45, 6.0, 2.55, accent=accent)
    L = add_textbox(s, 0.75, 4.60, 5.6, 2.40)
    add_para(L.text_frame, "Trigonometric phase encoding",
              size=13, bold=True, color=accent)
    add_para(L.text_frame, "𝒯(φ) = (sin φ, cos φ)",
              size=14, bold=True, color=LOW_BAND, space_before=4)
    add_para(L.text_frame, "= (b / (M+ε),  a / (M+ε))",
              size=11.5, color=MUTED_INK, space_before=2)
    add_para(L.text_frame, "▸  ±π 경계에서 continuous · gradient stable",
              size=11, color=MUTED_INK, space_before=6)
    add_para(L.text_frame, "▸  ε = 1e-8 — flat region에서 M → 0 안정화",
              size=11, color=MUTED_INK)

    stripe_card(s, 6.85, 4.45, 6.0, 2.55, accent=accent)
    R = add_textbox(s, 7.05, 4.60, 5.6, 2.40)
    add_para(R.text_frame, "Frequency Encoder",
              size=13, bold=True, color=accent)
    add_para(R.text_frame, "6.22 M params",
              size=24, bold=True, color=INK,
              align=PP_ALIGN.CENTER, space_before=8)
    add_para(R.text_frame, "≈ 3 % of the 208 M ControlNet",
              size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER,
              space_before=2)
    add_para(R.text_frame,
              "Identity-preserving init — γ ≡ 1, β ≡ 0 at step 0.",
              size=11, color=MUTED_INK, align=PP_ALIGN.CENTER,
              space_before=10)

    set_notes(s, (
        "SubbandBlock의 내부 구조 입니다. 하나의 DT-CWT scale을 입력으로 "
        "받아 (γ, β) modulation parameter를 생성합니다.\n\n"
        "위쪽 그림 — 처리 과정은 세 단계입니다. 첫째, complex coefficient "
        "C = a + ib에서 magnitude와 phase를 추출합니다. Magnitude branch "
        "에는 18 채널 (6 방향 × 3 RGB)이 들어가고, phase는 sin과 cos "
        "pair로 인코딩된 36 채널이 phase branch에 들어갑니다. 두 branch는 "
        "각각 DWConv → SiLU → Spatial Attention → Conv-1×1 의 동일한 "
        "operator를 통과합니다.\n\n"
        "둘째, recombination 단계 — magnitude를 cos과 sin과 element-wise "
        "곱해 complex 임베딩의 real부와 imaginary부를 만듭니다. "
        "이는 polar-to-Cartesian 관계 R·e^{iφ} = R·cos φ + i·R·sin φ 를 "
        "embedding space에서 그대로 반영한 것입니다.\n\n"
        "왼쪽 아래 — phase는 그대로 처리하면 ±π에서 wraparound 불연속이 "
        "생깁니다. φ = π - δ 와 φ = -π + δ 는 거의 같은 각도인데 raw "
        "값으로는 2π 차이가 납니다. 이 불연속은 convolution layer에 "
        "치명적입니다. 그래서 sin, cos pair로 인코딩해 unit circle 위에 "
        "올리면 continuous한 표현이 됩니다. 또 magnitude가 0에 가까운 "
        "flat region에서는 ε = 1e-8을 더해 수치 안정성을 보장합니다.\n\n"
        "오른쪽 아래 — 전체 Frequency Encoder는 6.22 M 파라미터로 "
        "ControlNet의 약 3% 입니다. Identity-preserving initialization — "
        "두 번째 conv의 weight를 0으로, γ의 bias를 1로, β의 bias를 0으로 "
        "초기화 — 덕분에 학습 시작 시점에 wrapped U-Net이 원래 U-Net과 "
        "동일하게 동작합니다."
    ))


def slide_bdsft_injection(prs):
    s, accent = add_content_slide(prs, section=3)
    set_section_title(s, 3, "Method · BD-SFT Injection",
                       subtitle="Asymmetric encoder-decoder modulation")

    # The U-Net injection figure
    add_picture_fit(s, FIG_DIR / "unet_injection.png",
                     0.55, 0.95, 12.3, 3.5)

    # Three rationale cards
    cards = [
        ("① Architectural alignment",
         "block_out_channels = [256, 512, 512, 1024].\n"
         "down_blocks[1] = up_blocks[1] = 512 ch — matches (γ, β) by "
         "construction."),
        ("② Functional separation",
         "FreeU (Si et al., 2024) — encoder suppresses HF & consolidates "
         "structure; decoder reintroduces HF via upsampling.\n"
         "→  align LOW ↔ encoder, HIGH ↔ decoder."),
        ("③ Conservative design",
         "Single injection per band — minimal overhead.\n"
         "Clean identity init at both injection points."),
    ]
    for i, (h, body) in enumerate(cards):
        x = 0.55 + i * 4.16
        stripe_card(s, x, 4.65, 4.0, 2.35, accent=accent)
        tb = add_textbox(s, x + 0.18, 4.78, 3.65, 2.20)
        add_para(tb.text_frame, h, size=12.5, bold=True, color=accent)
        for line in body.split("\n"):
            add_para(tb.text_frame, line, size=10.5, color=MUTED_INK,
                      space_before=3, space_after=2)

    set_notes(s, (
        "BD-SFT Injection scheme — band-decoupled spatial feature "
        "transform이라는 이름의 핵심 아이디어 입니다.\n\n"
        "위쪽 그림 — Stable Diffusion U-Net에 두 개의 injection point가 "
        "있습니다. 빨간색으로 표시된 down_blocks[1]은 encoder의 두 번째 "
        "block으로, LOW band SFT가 들어갑니다. 파란색으로 표시된 "
        "up_blocks[1]은 decoder의 두 번째 block으로, HIGH band SFT가 "
        "들어갑니다.\n\n"
        "왜 이렇게 asymmetric하게 배치했는지 세 가지 근거가 있습니다.\n\n"
        "첫째, architectural alignment — Stable Diffusion U-Net의 "
        "block_out_channels는 [256, 512, 512, 1024] 입니다. "
        "down_blocks[1]과 up_blocks[1]이 둘 다 512 채널 이고, 이건 "
        "본 연구의 (γ, β) 출력 차원과 정확히 일치합니다. 다른 block을 "
        "쓰면 추가적인 channel projection이 필요합니다.\n\n"
        "둘째, functional separation — FreeU 연구 (Si et al. 2024)가 "
        "밝힌 바, diffusion U-Net의 encoder는 high-frequency를 "
        "suppress하고 coarse structure를 consolidate합니다. 반면 "
        "decoder는 upsampling과 skip-connection fusion을 통해 "
        "high-frequency detail을 재도입합니다. 본 연구의 injection은 "
        "이 functional role과 정확히 align됩니다 — LOW band는 encoder의 "
        "structural abstraction을 강화하고, HIGH band는 decoder의 "
        "detail synthesis를 강화합니다.\n\n"
        "셋째, conservative design — band 당 injection을 하나로 "
        "제한함으로써 parameter overhead를 최소화하고, 두 injection point "
        "모두에서 clean한 identity initialization을 유지할 수 있습니다."
    ))


def slide_training_loss(prs):
    s, accent = add_content_slide(prs, section=3)
    set_section_title(s, 3, "Method · Training Loss",
                       subtitle="Band-decoupled wavelet supervision")

    # Total loss
    add_card(s, 0.55, 1.0, 12.3, 1.0, fill=SOFT_BG2, line=LIGHT_GRAY,
              line_width=0.5)
    add_left_stripe(s, 0.55, 1.0, 1.0, accent, thickness=0.08)
    tb = add_textbox(s, 0.55, 1.10, 12.3, 0.85, anchor=MSO_ANCHOR.MIDDLE)
    add_para(tb.text_frame, "Total objective",
              size=10.5, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_para(tb.text_frame,
              "ℒ_total  =  ℒ_MSE  +  λ_wav · 𝟙[t_step mod K = 0] · ℒ_wav      "
              "(λ_wav = 1.0,  K = 4)",
              size=15, bold=True, color=INK, align=PP_ALIGN.CENTER,
              space_before=4)

    # Wavelet loss equations
    add_card(s, 0.55, 2.2, 12.3, 1.85, fill=WHITE, line=DEEP, line_width=0.8)
    add_accent_strip(s, 0.55, 2.2, 12.3, DEEP, thickness=0.05)
    tb = add_textbox(s, 0.55, 2.35, 12.3, 1.7, anchor=MSO_ANCHOR.MIDDLE)
    add_para(tb.text_frame, "Band-decoupled wavelet loss",
              size=11, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    add_para(tb.text_frame,
              "ℒ_wav  =  λ_ℋ Σ_{j∈ℋ} ‖|Ĉ⁽ʲ⁾| − |C⁽ʲ⁾|‖₁   +   "
              "λ_ℒ Σ_{j∈ℒ} ℒⱼ^{mp}   +   λ_LP ‖Î^{LP} − I^{LP}‖₁",
              size=13, bold=True, color=INK, align=PP_ALIGN.CENTER,
              space_before=4)
    add_para(tb.text_frame,
              "ℒⱼ^{mp}  =  ‖|Ĉ⁽ʲ⁾| − |C⁽ʲ⁾|‖₁  +  𝔼[ |C⁽ʲ⁾| · "
              "(1 − cos(φ̂ − φ)) ]",
              size=11.5, italic=True, color=MUTED_INK, align=PP_ALIGN.CENTER,
              space_before=2)

    # Asymmetric supervision
    stripe_card(s, 0.55, 4.25, 6.0, 2.7, accent=HIGH_BAND, stripe_left=True,
                 stripe_top=False)
    tb = add_textbox(s, 0.85, 4.40, 5.5, 2.55)
    add_para(tb.text_frame, "HIGH band — magnitude only",
              size=13, bold=True, color=HIGH_BAND)
    add_para(tb.text_frame,
              "Texture를 loosely supervise — network가 phase에 "
              "constraint 없이 자유롭게 detail을 합성하도록.",
              size=11, color=MUTED_INK, space_before=4)
    add_para(tb.text_frame, "λ_ℋ = 0.1  ·  soft texture prior",
              size=12, bold=True, color=INK, space_before=8)

    stripe_card(s, 6.85, 4.25, 6.0, 2.7, accent=LOW_BAND, stripe_left=True,
                 stripe_top=False)
    tb = add_textbox(s, 7.15, 4.40, 5.5, 2.55)
    add_para(tb.text_frame, "LOW band — magnitude-weighted angular distance",
              size=13, bold=True, color=LOW_BAND)
    add_para(tb.text_frame,
              "Structural phase 오차를 |C| 비례로 강력하게 supervise — "
              "high-energy 영역의 phase가 더 중요.",
              size=11, color=MUTED_INK, space_before=4)
    add_para(tb.text_frame, "λ_ℒ = λ_LP = 1.0  ·  strict match",
              size=12, bold=True, color=INK, space_before=8)

    set_notes(s, (
        "마지막 method 슬라이드 — training loss 입니다.\n\n"
        "전체 objective는 두 term의 합 입니다. 첫째는 latent space의 "
        "표준 ε-MSE — diffusion model 학습의 기본입니다. 둘째는 "
        "band-decoupled wavelet loss를 K=4 step마다 적용합니다. K=4로 "
        "한 이유는 wavelet loss를 계산하려면 VAE로 latent를 pixel space "
        "로 decode해야 하는데 이게 비용이 크기 때문에, 매 step이 아닌 "
        "주기적으로 적용하는 절충입니다.\n\n"
        "Wavelet loss는 세 부분의 합 입니다. HIGH band 항은 magnitude만 "
        "L1, LOW band 항은 magnitude L1과 phase term의 합인 ℒⱼ^mp, 그리고 "
        "lowpass 항은 LP component의 L1 입니다.\n\n"
        "Phase term — 1 - cos(Δφ) — 는 directional statistics에서 표준 "
        "angular distance이고 wraparound-safe합니다. Ground truth "
        "magnitude |C|로 가중치를 주는 이유는, structural energy가 높은 "
        "영역에서는 phase fidelity가 중요하고 flat 영역에서는 phase가 "
        "수학적으로 undefined에 가깝기 때문에 supervision 가중치를 "
        "낮추는 게 맞습니다.\n\n"
        "Asymmetric supervision은 의도된 설계 입니다. HIGH band는 texture "
        "이므로 magnitude만 강제하고 phase는 자유롭게 두어 network가 "
        "실감 나는 detail을 만들 수 있게 합니다. LOW band는 structure "
        "이므로 phase까지 supervise해 정확한 구조 매칭을 강제합니다.\n\n"
        "Loss weight는 λ_H = 0.1, λ_L = λ_LP = 1.0 입니다. λ_H = 1.0으로 "
        "두면 HF artifact가 보이는 경향이 있었고, 0에 가까우면 perceptual "
        "quality가 떨어졌습니다. 0.1이 안정적인 절충 값이었습니다."
    ))


# ══════════════════════════════════════════════════════════════════════════
# Section 4 — Experiments
# ══════════════════════════════════════════════════════════════════════════
def slide_experiments_setup(prs):
    s, accent = add_content_slide(prs, section=4)
    set_section_title(s, 4, "Datasets & Training Setup")

    # Top — datasets table
    rows = [
        ("REDS (train split)", "Training",           "236",       "720p"),
        ("REDS4",              "In-domain eval",     "4 (000/011/015/020)", "720p"),
        ("Vid4",               "OOD eval",           "4",         "SD"),
        ("UDM10",              "OOD eval",           "10",        "720p"),
        ("SPMCS",              "OOD eval",           "30",        "SD"),
    ]
    add_table(s, 0.55, 1.0, 7.4, 2.4,
              headers=["Dataset", "Role", "Seqs", "Resolution"],
              rows=rows, header_size=11, body_size=10.5)

    # Right — training config
    rows2 = [
        ("Backbone (frozen)", "SD v2.1 U-Net 472 M  +  VAE 55 M"),
        ("Trainable",         "ControlNet 208 M  +  Freq.Encoder 6.22 M"),
        ("Optimizer",         "AdamW · lr 1e-4 · 20k iters"),
        ("Batch / window",    "B=8 · T_win=3 · 64² → 256² patches"),
        ("DT-CWT",            "near_sym_a / qshift_a, J=4"),
        ("Inference",         "50 DDPM · Bidirectional Sampling"),
    ]
    add_table(s, 8.15, 1.0, 4.75, 3.5,
              headers=["Item", "Specification"],
              rows=rows2, header_size=10, body_size=9.5)

    # Bottom — controlled comparison callout
    add_card(s, 0.55, 3.65, 7.4, 3.3, fill=SOFT_BG2, line=LIGHT_GRAY,
              line_width=0.5)
    add_left_stripe(s, 0.55, 3.65, 3.3, accent, thickness=0.10)
    tb = add_textbox(s, 0.85, 3.83, 6.95, 3.10)
    add_para(tb.text_frame, "Controlled comparison",
              size=13, bold=True, color=accent)
    add_para(tb.text_frame,
              "▸  Training data · backbone · temporal scaffold — all matched "
              "to StableVSR.",
              size=11.5, color=MUTED_INK, space_before=6)
    add_para(tb.text_frame,
              "▸  유일한 architectural 차이 — Wavelet Conditioning Module "
              "(WCM) + BD-SFT injection.",
              size=11.5, color=MUTED_INK, space_before=4)
    add_para(tb.text_frame,
              "⇒  REDS4 performance gap이 WCM에 의해 isolation됨 — "
              "ablation-by-design.",
              size=12, bold=True, color=INK, space_before=10)

    set_notes(s, (
        "실험 셋업입니다.\n\n"
        "Dataset — 학습은 REDS train split에서 240개 시퀀스 중 REDS4의 "
        "4개 시퀀스를 제외한 236개로 진행합니다. 이건 BasicVSR과 "
        "StableVSR의 표준 컨벤션을 따른 것입니다. 평가는 in-domain "
        "REDS4와 out-of-distribution Vid4, UDM10, SPMCS 입니다.\n\n"
        "Training config — AdamW 옵티마이저, learning rate 1e-4 constant, "
        "20,000 iteration. Batch size 8, temporal window 3, 64×64 LR "
        "patch와 256×256 HR patch 입니다. NVIDIA RTX A6000 48GB 2장으로 "
        "학습했습니다.\n\n"
        "가장 중요한 점은 이 실험이 controlled comparison으로 설계됐다는 "
        "것입니다. Training data, backbone (SD v2.1), temporal scaffold "
        "(ControlNet + RAFT + bidirectional sampling) — 모두 StableVSR와 "
        "동일합니다. 유일한 차이는 Wavelet Conditioning Module과 BD-SFT "
        "injection의 추가 뿐입니다. 따라서 REDS4 결과의 성능 차이가 "
        "전적으로 WCM에 기인한다고 확실히 attribute할 수 있습니다 — "
        "ablation-by-design 입니다."
    ))


def slide_metrics(prs):
    s, accent = add_content_slide(prs, section=4)
    set_section_title(s, 4, "Evaluation Metrics",
                       subtitle="9 metrics across 4 categories")

    cats = [
        ("Reconstruction fidelity",
         "PSNR ↑   ·   SSIM ↑",
         "Pixel · structural similarity.",
         SEC[5]),
        ("Full-reference perceptual",
         "LPIPS ↓   ·   DISTS ↓",
         "Deep-feature perceptual distance.",
         SEC[2]),
        ("No-reference perceptual",
         "MUSIQ ↑   ·   CLIP-IQA ↑   ·   NIQE ↓",
         "Reference-free quality — critical for generative VSR.",
         SEC[3]),
        ("Temporal consistency",
         "tLPIPS ↓   ·   tOF ↓",
         "Inter-frame perceptual / motion stability.",
         SEC[1]),
    ]
    for i, (cat, metrics, body, col) in enumerate(cats):
        row = i // 2
        cidx = i % 2
        x = 0.55 + cidx * 6.15
        y = 1.0 + row * 2.0
        stripe_card(s, x, y, 6.0, 1.85, accent=col, stripe_left=True,
                     stripe_top=False)
        tb = add_textbox(s, x + 0.25, y + 0.13, 5.65, 1.65)
        add_para(tb.text_frame, cat, size=13, bold=True, color=col)
        add_para(tb.text_frame, metrics, size=14, bold=True, color=INK,
                  space_before=4)
        add_para(tb.text_frame, body, size=11, color=MUTED_INK,
                  space_before=2)

    # Bottom — perception-distortion footnote
    add_card(s, 0.55, 5.10, 12.3, 1.85, fill=SOFT_BG2, line=LIGHT_GRAY,
              line_width=0.5)
    add_left_stripe(s, 0.55, 5.10, 1.85, accent, thickness=0.10)
    tb = add_textbox(s, 0.85, 5.25, 11.85, 1.65)
    add_para(tb.text_frame, "Why so many metrics?",
              size=12.5, bold=True, color=accent)
    add_para(tb.text_frame,
              "Generative VSR은 perception–distortion trade-off 위에 "
              "위치 — 단일 metric으로 전체 그림을 보여줄 수 없습니다. "
              "Pixel metric은 사실적인 합성을 페널라이즈하고, NR perceptual "
              "metric은 reference-based metric을 보완하며, temporal metric은 "
              "per-frame metric이 놓치는 flicker를 포착합니다.",
              size=11.5, color=MUTED_INK, space_before=4)

    set_notes(s, (
        "총 9개의 평가 지표를 4개 카테고리로 분류합니다.\n\n"
        "Reconstruction fidelity — PSNR과 SSIM. Pixel-level 또는 "
        "structural similarity로, generative model에는 다소 불리한 "
        "지표입니다.\n\n"
        "Full-reference perceptual — LPIPS와 DISTS. Deep feature 기반 "
        "perceptual distance입니다. DISTS는 texture variation에 robust해 "
        "generative VSR 평가에 특히 적합합니다.\n\n"
        "No-reference perceptual — MUSIQ, CLIP-IQA, NIQE. Reference 없이 "
        "quality를 평가하므로 generative model의 사실적 합성을 "
        "underestimate하지 않습니다. 세 가지가 각각 transformer 기반, "
        "vision-language 기반, natural scene statistics 기반이라 서로를 "
        "보완합니다.\n\n"
        "Temporal consistency — tLPIPS와 tOF. tLPIPS는 인접 프레임 간의 "
        "perceptual difference로 texture-level flicker를 잘 잡아냅니다. "
        "diffusion 기반 VSR 평가에 특히 중요합니다. tOF는 optical flow "
        "차이로 motion fidelity를 측정합니다.\n\n"
        "이렇게 다양한 metric을 쓰는 이유는 — generative VSR이 "
        "perception-distortion trade-off 위에 있어서 단일 metric으로는 "
        "공정한 평가가 불가능하기 때문입니다."
    ))


# ══════════════════════════════════════════════════════════════════════════
# Section 5 — Results
# ══════════════════════════════════════════════════════════════════════════
def slide_results_reds4(prs):
    s, accent = add_content_slide(prs, section=5)
    set_section_title(s, 5, "Results · REDS4 — Direct vs. StableVSR")

    rows = [
        ("StableVSR",
         "24.04", "0.690", "0.309", "0.164", "42.79", "0.237", "4.38",
         "41.11", "13.962"),
        ("WC-BD-SFT (Ours)",
         "24.48", "0.691", "0.193", "0.088", "65.70", "0.386", "2.77",
         "15.25", "11.118"),
        ("Δ (relative)",
         "+1.8%", "+0.1%", "−37.5%", "−46.3%", "+53.5%", "+62.9%",
         "−36.8%", "−62.9%", "−20.4%"),
    ]
    add_table(s, 0.4, 1.0, 12.55, 1.8,
              headers=["Method",
                       "PSNR↑", "SSIM↑", "LPIPS↓", "DISTS↓",
                       "MUSIQ↑", "CLIP-IQA↑", "NIQE↓", "tLPIPS↓", "tOF↓"],
              rows=rows, header_size=10, body_size=10,
              highlight_row_idx=1, highlight_fill=RGBColor(0xE3, 0xF0, 0xFF))

    # KPI cards
    big_stat(s, 0.55, 3.1, 4.0, 1.8, label="Perceptual · LPIPS ↓",
              value="−37.5 %", sub="0.309 → 0.193", accent=SEC[2])
    big_stat(s, 4.7, 3.1, 4.0, 1.8, label="Perceptual · MUSIQ ↑",
              value="+53.5 %", sub="42.79 → 65.70", accent=SEC[3])
    big_stat(s, 8.85, 3.1, 4.0, 1.8, label="Temporal · tLPIPS ↓",
              value="−62.9 %", sub="41.11 → 15.25", accent=accent)

    # Interpretation
    add_card(s, 0.55, 5.1, 12.3, 1.85, fill=SOFT_BG2, line=LIGHT_GRAY,
              line_width=0.5)
    add_left_stripe(s, 0.55, 5.1, 1.85, accent, thickness=0.10)
    tb = add_textbox(s, 0.85, 5.25, 11.85, 1.6)
    add_para(tb.text_frame, "Interpretation",
              size=13, bold=True, color=accent)
    add_para(tb.text_frame,
              "Training data · backbone · temporal scaffold이 모두 "
              "matched된 상태이므로, 이 성능 gap은 wavelet-conditioned "
              "BD-SFT mechanism의 단독 기여로 정확하게 attribute됩니다.",
              size=12, color=MUTED_INK, space_before=4)

    set_notes(s, (
        "이 슬라이드가 본 논문에서 가장 직접적인 비교 입니다.\n\n"
        "StableVSR 대비 9개 metric 전부에서 향상되었습니다. PSNR도 "
        "0.44 dB 올라가고, 그 외 perceptual / temporal metric은 모두 "
        "큰 폭으로 개선됐습니다.\n\n"
        "특히 주목할 세 가지 — LPIPS는 0.309에서 0.193으로 37.5% 감소, "
        "MUSIQ는 42.79에서 65.70으로 53.5% 증가, tLPIPS는 41.11에서 "
        "15.25로 무려 62.9% 감소입니다. 3D attention이나 full-network "
        "fine-tuning 없이 이 정도 temporal consistency 향상을 달성한 "
        "것입니다.\n\n"
        "Training data, backbone, temporal scaffold가 모두 동일하므로 — "
        "유일한 차이는 WCM의 유무 입니다. 따라서 이 gap은 정확히 본 "
        "연구의 기여라고 attribute할 수 있습니다. Decoder의 HIGH band "
        "injection 덕분에 texture synthesis 단계에서 강한 conditioning이 "
        "들어가, perceptual axis에서의 향상이 가장 두드러집니다."
    ))


def slide_results_dm_group(prs):
    s, accent = add_content_slide(prs, section=5)
    set_section_title(s, 5, "Results · REDS4 — Within DM Group",
                       subtitle="vs. DGAF-VSR · BasicVSR++ cross-paradigm")

    rows = [
        ("non-DM", "BasicVSR++",
         "24.88", "0.730", "0.364", "0.179", "41.23", "0.265", "5.90",
         "38.26", "13.739", "—"),
        ("DM", "StableVSR",
         "24.04", "0.690", "0.309", "0.164", "42.79", "0.237", "4.38",
         "41.11", "13.962", "3.00"),
        ("DM", "DGAF-VSR",
         "24.07", "0.694", "0.307", "0.161", "43.41", "0.242", "4.37",
         "40.12", "13.732", "1.89"),
        ("DM", "WC-BD-SFT (Ours)",
         "24.48", "0.691", "0.193", "0.088", "65.70", "0.386", "2.77",
         "15.25", "11.118", "1.11"),
    ]
    add_table(s, 0.3, 1.0, 12.75, 2.4,
              headers=["Paradigm", "Method",
                       "PSNR↑", "SSIM↑", "LPIPS↓", "DISTS↓",
                       "MUSIQ↑", "CLIP-IQA↑", "NIQE↓",
                       "tLPIPS↓", "tOF↓", "Mean rank"],
              rows=rows, header_size=9, body_size=8.5,
              highlight_row_idx=3, highlight_fill=RGBColor(0xE3, 0xF0, 0xFF))

    # Two takeaways
    stripe_card(s, 0.55, 3.7, 6.0, 3.3, accent=SEC[3])
    tb = add_textbox(s, 0.75, 3.85, 5.6, 3.15)
    add_para(tb.text_frame, "vs. DGAF-VSR (strongest DM baseline)",
              size=13, bold=True, color=SEC[3])
    for line, gain in [
        ("LPIPS  0.307 → 0.193", "−37.1 %"),
        ("MUSIQ  43.41 → 65.70", "+51.4 %"),
        ("tLPIPS 40.12 → 15.25", "−62.0 %"),
    ]:
        p = tb.text_frame.add_paragraph()
        p.space_before = Pt(4)
        r1 = p.add_run()
        r1.text = "▸  " + line + "    "
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = MUTED_INK
        r2 = p.add_run()
        r2.text = gain
        r2.font.size = Pt(11.5)
        r2.font.bold = True
        r2.font.color.rgb = SEC[3]
    add_para(tb.text_frame, "Mean rank within DM:  1.11  (vs. 1.89 / 3.00)",
              size=12, bold=True, color=INK, space_before=10)

    stripe_card(s, 6.85, 3.7, 6.0, 3.3, accent=DEEP)
    tb = add_textbox(s, 7.05, 3.85, 5.6, 3.15)
    add_para(tb.text_frame, "Cross-paradigm reference — BasicVSR++",
              size=13, bold=True, color=DEEP)
    add_para(tb.text_frame, "PSNR / SSIM은 BasicVSR++가 우수 (regression).",
              size=11.5, color=MUTED_INK, space_before=4)
    add_para(tb.text_frame,
              "Perceptual 4 metric에서는 모두 열위 — Blau & Michaeli "
              "(2018)의 perception-distortion trade-off.",
              size=11.5, color=MUTED_INK, space_before=4)
    add_para(tb.text_frame,
              "DM 내부 ranking이 본 연구 mechanism의 공정한 비교 축.",
              size=11.5, italic=True, color=GRAY, space_before=6)

    set_notes(s, (
        "REDS4 in-domain 비교를 좀 더 넓은 범위로 확장한 결과 입니다.\n\n"
        "Diffusion model 그룹 내에서 본 연구의 WC-BD-SFT는 9개 metric "
        "중 8개에서 best, 1개(SSIM)에서 second-best입니다. Mean rank가 "
        "1.11으로 DGAF-VSR의 1.89, StableVSR의 3.00에 비해 명확히 "
        "앞섭니다.\n\n"
        "왼쪽 — DGAF-VSR과의 비교 입니다. DGAF-VSR는 CVPR 2026에 "
        "발표된 가장 최근의 diffusion-based VSR로, 본 연구의 가장 강한 "
        "baseline 입니다. LPIPS 37.1% 감소, MUSIQ 51.4% 증가, tLPIPS "
        "62.0% 감소로 perceptual axis에서 일관된 큰 폭의 향상이 "
        "있습니다.\n\n"
        "오른쪽 — BasicVSR++는 cross-paradigm reference 입니다. CNN 기반 "
        "regression model이라 PSNR과 SSIM에서는 우수하지만, 4개 "
        "perceptual metric 모두에서 열위 입니다. 이건 Blau와 Michaeli의 "
        "perception-distortion trade-off의 전형적인 패턴 입니다. "
        "다른 paradigm을 직접 ranking으로 비교하는 건 공정하지 않으므로, "
        "DM 내부 ranking을 본 연구 mechanism의 공정한 비교 축으로 "
        "삼습니다."
    ))


def slide_results_cross_dataset(prs):
    s, accent = add_content_slide(prs, section=5)
    set_section_title(s, 5, "Results · Cross-Dataset",
                       subtitle="Vid4 · UDM10 · SPMCS — OOD evaluation")

    rows = [
        ("Vid4",  "BasicVSR++",  "26.26", "0.828", "0.189", "61.50", "0.341", "5.04", "15.12"),
        ("Vid4",  "StableVSR",   "22.98", "0.674", "0.185", "67.22", "0.454", "3.19", "25.36"),
        ("Vid4",  "DGAF-VSR",    "23.29", "0.690", "0.177", "67.96", "0.470", "3.10", "17.39"),
        ("Vid4",  "Ours",        "22.64", "0.664", "0.194", "66.15", "0.408", "3.32", "28.53"),
        ("UDM10", "BasicVSR++",  "37.48", "0.956", "0.060", "59.36", "0.443", "5.60",  "5.55"),
        ("UDM10", "StableVSR",   "26.71", "0.834", "0.100", "55.69", "0.362", "4.66",  "4.42"),
        ("UDM10", "DGAF-VSR",    "26.71", "0.835", "0.099", "57.15", "0.380", "4.61",  "2.99"),
        ("UDM10", "Ours",        "25.54", "0.811", "0.124", "63.20", "0.447", "4.12", "14.48"),
        ("SPMCS", "BasicVSR++",  "21.94", "0.617", "0.187", "62.48", "0.434", "5.17",  "3.60"),
        ("SPMCS", "StableVSR",   "19.42", "0.478", "0.196", "69.98", "0.582", "3.28", "51.11"),
        ("SPMCS", "DGAF-VSR",    "20.06", "0.511", "0.178", "67.70", "0.533", "3.61", "20.10"),
        ("SPMCS", "Ours",        "19.94", "0.506", "0.183", "67.22", "0.503", "3.70", "31.59"),
    ]
    add_table(s, 0.35, 1.0, 7.3, 5.95,
              headers=["Set", "Method",
                       "PSNR↑", "SSIM↑", "LPIPS↓", "MUSIQ↑",
                       "CLIP-IQA↑", "NIQE↓", "tLPIPS↓"],
              rows=rows, header_size=9, body_size=8.5)

    # Right column — observations
    obs = [
        ("UDM10 — NR perceptual wins",
         "MUSIQ +13.5 % · CLIP-IQA +23.5 % · NIQE −11.6 % (in DM group).\n"
         "Higher tLPIPS — short static clips, see Discussion.",
         SEC[3]),
        ("SPMCS — temporal gain",
         "tLPIPS 51.11 → 31.59 (−38.2 %) over StableVSR.\n"
         "Long-motion content에서는 frequency conditioning이 "
         "stability에 도움.",
         SEC[2]),
        ("Vid4 — distribution gap",
         "Compressed SD vs. clean REDS 720p — gain 감소.\n"
         "DGAF의 OGWM과 fusion이 future work.",
         DEEP),
    ]
    for i, (h, body, col) in enumerate(obs):
        y = 1.0 + i * 1.95
        stripe_card(s, 7.95, y, 4.95, 1.85, accent=col)
        tb = add_textbox(s, 8.12, y + 0.13, 4.65, 1.7)
        add_para(tb.text_frame, h, size=12, bold=True, color=col)
        add_para(tb.text_frame, body, size=10.5, color=MUTED_INK,
                  space_before=4)

    set_notes(s, (
        "Out-of-distribution 평가 결과 입니다. 흥미로운 패턴이 보입니다.\n\n"
        "UDM10에서는 DM 그룹 내 NR perceptual metric을 모두 1위로 받았습니다 "
        "— MUSIQ 13.5%, CLIP-IQA 23.5%, NIQE 11.6% 향상. 다만 tLPIPS는 "
        "14.48로 다른 방법보다 높습니다. UDM10은 32 프레임의 짧고 거의 "
        "정적인 시퀀스 이고, wavelet-driven texture synthesis가 만든 "
        "프레임별 미세한 차이가 motion으로 흡수되지 못해 flicker로 "
        "축적된 결과 입니다. 이건 Discussion에서 자세히 다룹니다.\n\n"
        "SPMCS는 반대로 long-motion content 입니다. tLPIPS가 51.11에서 "
        "31.59로 38.2% 감소 — frequency conditioning이 충분한 motion이 "
        "있을 때는 temporal stability에 도움이 됩니다.\n\n"
        "Vid4는 compressed SD 콘텐츠로 REDS의 clean 720p와 distribution "
        "차이가 가장 큽니다. 본 연구의 wavelet prior가 REDS-like statistic "
        "에 specialize되어 있어 gain이 작습니다. DGAF의 high-resolution "
        "feature warping이 이쪽에서 강점을 보이는데, 이건 본 연구와 "
        "orthogonal하므로 fusion이 자연스러운 future work 입니다."
    ))


def slide_ablation(prs):
    s, accent = add_content_slide(prs, section=5)
    set_section_title(s, 5, "Ablation · Frequency-Band Injection",
                       subtitle="Inference-time band disabling on REDS4")

    rows = [
        ("WC-BD-SFT (full)",     "24.48", "0.691", "0.193", "0.088",
         "65.70", "0.386", "2.77", "15.25", "11.118"),
        ("w/o HIGH",             "24.34", "0.689", "0.223", "0.104",
         "64.49", "0.365", "3.00", "22.11", "11.984"),
        ("w/o LOW",              "24.70", "0.701", "0.210", "0.098",
         "63.39", "0.335", "3.00", "16.57", "11.392"),
        ("w/o both ≈ StableVSR", "24.46", "0.695", "0.247", "0.119",
         "59.35", "0.298", "3.29", "25.18", "12.278"),
    ]
    add_table(s, 0.3, 1.0, 12.75, 2.3,
              headers=["Variant", "PSNR↑", "SSIM↑", "LPIPS↓", "DISTS↓",
                       "MUSIQ↑", "CLIP-IQA↑", "NIQE↓",
                       "tLPIPS↓", "tOF↓"],
              rows=rows, header_size=10, body_size=10,
              highlight_row_idx=0,
              highlight_fill=RGBColor(0xE3, 0xF0, 0xFF))

    insights = [
        ("HIGH is critical",
         "tLPIPS  15.25 → 22.11   (+45.0 %)\n"
         "LPIPS   0.193 → 0.223   (+15.5 %)",
         HIGH_BAND),
        ("LOW shapes perception",
         "MUSIQ    65.70 → 63.39\n"
         "CLIP-IQA 0.386 → 0.335\n"
         "(PSNR / SSIM slightly rise — trade-off)",
         LOW_BAND),
        ("Super-additive when both off",
         "LPIPS = 0.247 — worse than either single disable.\n"
         "MUSIQ collapses to 59.35.\n"
         "⇒ bands are jointly necessary.",
         DEEP),
    ]
    for i, (h, body, col) in enumerate(insights):
        x = 0.55 + i * 4.16
        stripe_card(s, x, 3.55, 4.0, 3.4, accent=col)
        tb = add_textbox(s, x + 0.18, 3.70, 3.65, 3.2)
        add_para(tb.text_frame, h, size=13, bold=True, color=col)
        add_para(tb.text_frame, body, size=11.5, color=MUTED_INK,
                  space_before=6)

    set_notes(s, (
        "Band injection ablation 입니다. 학습된 동일 모델에서 inference "
        "시점에 SFT modulation을 선택적으로 disable (γ=1, β=0)해 각 "
        "band의 기여를 isolate 했습니다. 별도 학습이 아니라 inference "
        "시점 ablation이라 학습 trajectory의 차이로 인한 confound가 "
        "없습니다.\n\n"
        "HIGH를 disable하면 모든 perceptual / temporal metric이 크게 "
        "악화 됩니다. tLPIPS는 15.25에서 22.11로 45% 증가, LPIPS는 "
        "15.5% 증가 합니다. Decoder에서 texture detail을 합성하는 "
        "stage의 HIGH-band guidance가 inter-frame texture coherence에 "
        "결정적이라는 게 확인됩니다.\n\n"
        "LOW를 disable하면 패턴이 다릅니다 — PSNR과 SSIM은 미세하게 "
        "오히려 좋아지지만, NR perceptual metric은 떨어집니다. Encoder "
        "에서의 low-band guidance가 structural representation을 perceptually "
        "natural한 방향으로 shaping 하는데, 그게 disable되면 fidelity "
        "쪽으로 살짝 회귀하는 거죠. 이건 perception-distortion trade-off "
        "의 미니어처 입니다.\n\n"
        "가장 중요한 발견은 super-additive degradation 입니다. 둘 다 "
        "disable하면 LPIPS가 0.247로, 단일 disable의 worst case인 0.223 "
        "보다 더 나빠집니다. MUSIQ도 59.35로 폭락합니다. 즉 두 band가 "
        "독립적으로 좋은 게 아니라, U-Net 내부에서 비선형적으로 "
        "interaction하며 jointly necessary 하다는 의미입니다. Asymmetric "
        "injection이 단순한 additive 개선이 아니라 synergistic이라는 "
        "증거입니다."
    ))


def slide_discussion(prs):
    s, accent = add_content_slide(prs, section=5)
    set_section_title(s, 5, "Discussion · Trade-offs Revealed")

    # Left — P-D curve diagram
    add_picture_fit(s, FIG_DIR / "pd_curve.png",
                     0.55, 1.0, 6.0, 5.95)

    # Right — three trade-off cards
    cards = [
        ("① Perception–distortion (paradigms)",
         "BasicVSR++ → high pixel fidelity.\n"
         "DM 방법 → high perceptual quality.\n"
         "WC-BD-SFT는 perceptual 극단 — REDS4 / UDM10 NR 1위.",
         SEC[5]),
        ("② Specialization–generalization",
         "REDS-only 학습 → WCM 기여 isolation.\n"
         "REDS-like statistics에 강함, Vid4 SD에서는 gain 감소.",
         SEC[3]),
        ("③ Temporal stability ↔ motion",
         "UDM10 (32-frame, 정적) → tLPIPS 악화.\n"
         "SPMCS (long-motion) → tLPIPS 개선.",
         SEC[2]),
    ]
    for i, (h, body, col) in enumerate(cards):
        y = 1.0 + i * 2.0
        stripe_card(s, 6.85, y, 6.0, 1.85, accent=col, stripe_left=True,
                     stripe_top=False)
        tb = add_textbox(s, 7.15, y + 0.13, 5.65, 1.65)
        add_para(tb.text_frame, h, size=12, bold=True, color=col)
        for line in body.split("\n"):
            add_para(tb.text_frame, line, size=10.5, color=MUTED_INK,
                      space_before=2, space_after=2)

    set_notes(s, (
        "실험으로 드러난 세 가지 trade-off 입니다.\n\n"
        "첫째 — paradigm 간 perception-distortion trade-off 입니다. "
        "BasicVSR++는 CNN regression model이라 pixel fidelity를 최대화 "
        "하지만 perceptually over-smoothed 합니다. DM 방법들은 사실적 "
        "texture를 만들지만 pixel accuracy를 일부 희생합니다. 본 연구의 "
        "WC-BD-SFT는 perceptual 극단을 추구해서 REDS4와 UDM10에서 NR "
        "perceptual metric 1위 입니다.\n\n"
        "둘째 — specialization-generalization trade-off 입니다. 본 연구는 "
        "WCM 기여를 isolate하기 위해 의도적으로 REDS-only 학습을 "
        "선택했습니다. 그 결과 wavelet prior가 REDS의 frequency statistics "
        "에 specialize되어 — REDS4에서는 최강이지만 Vid4 같은 compressed "
        "SD에서는 gain이 작습니다. 이건 single-distribution training의 "
        "본질적 trade-off이지, mechanism의 한계가 아닙니다. 별도의 연구 "
        "질문 입니다.\n\n"
        "셋째 — sequence motion에 따른 temporal stability 입니다. UDM10 "
        "처럼 짧고 정적인 시퀀스에서는 wavelet-driven texture의 미세 "
        "variation이 flicker로 축적됩니다. 반면 SPMCS처럼 motion이 풍부한 "
        "long sequence에서는 tLPIPS가 향상됩니다. Bidirectional sampling이 "
        "motion을 stochastic variation을 흡수하는 매개로 활용한다는 "
        "가설로 설명 가능합니다."
    ))


# ══════════════════════════════════════════════════════════════════════════
# Section 6 — Conclusion
# ══════════════════════════════════════════════════════════════════════════
def slide_limitations(prs):
    s, accent = add_content_slide(prs, section=6)
    set_section_title(s, 6, "Limitations")

    items = [
        ("Training-distribution dependence",
         "REDS-like content에서 최강 · Vid4 SD에서 gain 감소.",
         "광범위한 학습 mix · degradation-aware augmentation"),
        ("Perception–distortion trade-off",
         "PSNR/SSIM을 perceptual quality에 일부 양보 — 특히 UDM10.",
         "Perceptual application에 적합 · 의료/과학 영상엔 부적합"),
        ("Short-static flicker",
         "UDM10 (32 frame)에서 DGAF/StableVSR보다 tLPIPS 높음.",
         "Lightweight temporal regularization (flow-guided latent warping)"),
        ("Computational overhead",
         "4-level DT-CWT가 frame당 modest cost 추가.",
         "Diffusion denoising cost에 비해 작음 · 추가 최적화 가능"),
        ("Comparison scope",
         "Single fixed seed · 2025년 다수 method 미평가.",
         "Multi-seed evaluation · RVRT / MGLD / DiffVSR 후속 비교"),
    ]
    for i, (h, situation, fix) in enumerate(items):
        col = i % 2
        row = i // 2
        x = 0.55 + col * 6.15
        y = 1.0 + row * 2.0
        if row == 2 and col == 1:
            continue
        stripe_card(s, x, y, 6.0, 1.85, accent=accent)
        tb = add_textbox(s, x + 0.2, y + 0.13, 5.65, 1.65)
        add_para(tb.text_frame, h, size=12.5, bold=True, color=accent)
        add_para(tb.text_frame, situation, size=11, color=MUTED_INK,
                  space_before=4)
        add_para(tb.text_frame, "→ " + fix, size=10.5, italic=True,
                  color=GRAY, space_before=2)

    set_notes(s, (
        "본 연구의 한계를 다섯 가지로 명시적으로 인정합니다.\n\n"
        "첫째 — training distribution 의존성. REDS-like content에서는 "
        "강하지만 Vid4 같은 compressed SD에서는 gain이 작습니다. 광범위한 "
        "data mixture와 degradation-aware augmentation으로 mitigation "
        "가능합니다.\n\n"
        "둘째 — perception-distortion trade-off가 실재합니다. PSNR과 "
        "SSIM을 perceptual quality에 일부 양보합니다. UDM10에서 PSNR이 "
        "BasicVSR++의 37.48 dB 대비 25.54 dB 입니다. 본 framework는 "
        "perceptual application에 적합하고, 의료 영상이나 과학 영상처럼 "
        "strict pixel accuracy가 필요한 영역에는 부적합 합니다.\n\n"
        "셋째 — short / quasi-static sequence에서 flicker가 더 심합니다. "
        "Lightweight temporal regularization 추가로 해결 가능합니다.\n\n"
        "넷째 — 4-level DT-CWT의 computational overhead. 다행히 diffusion "
        "denoising 비용에 비해 작아서 critical bottleneck은 아닙니다.\n\n"
        "다섯째 — comparison scope의 제한. 재현성을 위해 단일 seed(42)를 "
        "사용했고 multi-seed evaluation이 future work 입니다. 또한 "
        "release timing과 GPU 제약으로 RVRT, MGLD-VSR, DiffVSR, STAR, "
        "DC-VSR, DLoRAL, UltraVSR, SeedVR2 등 일부 최신 method는 "
        "벤치마크하지 못했습니다."
    ))


def slide_conclusion(prs):
    s, accent = add_content_slide(prs, section=6)
    set_section_title(s, 6, "Conclusion · Key Findings")

    # Summary card
    add_card(s, 0.55, 1.0, 12.3, 1.8, fill=SOFT_BG2, line=LIGHT_GRAY,
              line_width=0.5)
    add_left_stripe(s, 0.55, 1.0, 1.8, accent, thickness=0.10)
    tb = add_textbox(s, 0.85, 1.18, 11.85, 1.6)
    add_para(tb.text_frame,
              "Frequency-domain conditioning is an effective inductive bias "
              "for diffusion-based VSR.",
              size=15, bold=True, color=INK)
    add_para(tb.text_frame,
              "WC-BD-SFT injects DT-CWT priors into a frozen pre-trained "
              "diffusion U-Net via asymmetric encoder-decoder modulation, "
              "with magnitude-phase decoupling and a band-decoupled wavelet "
              "loss — at only 6.22 M added trainable parameters.",
              size=11.5, color=MUTED_INK, space_before=6)

    # Three KPI cards
    big_stat(s, 0.55, 3.0, 4.0, 4.0,
              label="REDS4 · tLPIPS",
              value="−62.9 %",
              sub="41.11  →  15.25\n\nvs. StableVSR baseline",
              accent=SEC[2])
    big_stat(s, 4.7, 3.0, 4.0, 4.0,
              label="REDS4 · HF spectral power",
              value="0.693",
              sub="vs. StableVSR's 0.341\n\n+103 % HF retention",
              accent=SEC[3])
    big_stat(s, 8.85, 3.0, 4.0, 4.0,
              label="Frequency Encoder",
              value="6.22 M",
              sub="≈ 3 % of ControlNet\n\nU-Net + VAE remain frozen",
              accent=SEC[5])

    set_notes(s, (
        "결론을 요약합니다.\n\n"
        "본 학위논문의 핵심 발견은 — frequency-domain conditioning이 "
        "diffusion-based VSR을 위한 effective하고 parameter-efficient한 "
        "inductive bias라는 것입니다. WC-BD-SFT는 DT-CWT prior를 frozen "
        "pre-trained diffusion U-Net에 asymmetric encoder-decoder injection "
        "으로 주입하며, magnitude-phase 분리, band-decoupled wavelet "
        "loss를 결합합니다. 추가되는 학습 파라미터는 6.22 M 입니다.\n\n"
        "세 개의 headline 숫자로 본 연구의 기여를 요약하겠습니다.\n\n"
        "첫째 — temporal consistency. In-domain REDS4에서 tLPIPS가 62.9% "
        "감소 했습니다.\n\n"
        "둘째 — frequency-domain fidelity. Ground truth 대비 high-frequency "
        "spectral power retention이 0.693으로, StableVSR의 0.341 대비 두 "
        "배 이상입니다.\n\n"
        "셋째 — parameter efficiency. Frequency Encoder는 6.22 M으로 "
        "ControlNet의 약 3%에 불과하고, 사전학습된 U-Net과 VAE는 학습 "
        "내내 frozen으로 유지되므로 generative prior가 그대로 보존 됩니다."
    ))


def slide_future_work(prs):
    s, accent = add_content_slide(prs, section=6)
    set_section_title(s, 6, "Future Work")

    items = [
        (1, "Flow-free pipeline",
         "ControlNet의 RAFT warping을 frequency-aware temporal module로 "
         "교체 — optical-flow 의존성 완전 제거."),
        (2, "Combine with DGAF-VSR",
         "Frequency-domain conditioning ⟂ feature-domain warping — 두 "
         "방향이 orthogonal해 fusion 시 cross-dataset 개선 기대."),
        (3, "Broader training mixture",
         "REDS + Vimeo-90K 학습 · degradation-aware augmentation — Vid4 "
         "같은 OOD에서의 gap 완화."),
        (4, "Latent-space wavelet analysis",
         "Wavelet loss를 latent space에서 수행 · adaptive band partitioning "
         "— VAE-decoding overhead 감소."),
        (5, "Short-sequence temporal reg.",
         "Flow-guided latent warping / multi-frame attention 같은 "
         "경량 module 추가 — UDM10 flicker 해결."),
    ]
    for i, (n, h, body) in enumerate(items):
        if i < 3:
            x = 0.55 + i * 4.16
            y = 1.0
        else:
            x = 2.65 + (i - 3) * 4.16
            y = 4.0
        w, ht = 4.0, 2.85
        stripe_card(s, x, y, w, ht, accent=accent)
        tb = add_textbox(s, x + 0.2, y + 0.18, w - 0.4, ht - 0.36)
        add_para(tb.text_frame, f"{n}.  {h}", size=13, bold=True, color=accent)
        add_para(tb.text_frame, body, size=11, color=MUTED_INK,
                  space_before=8)

    set_notes(s, (
        "마지막으로 future work 다섯 가지를 제시합니다.\n\n"
        "첫째 — flow-free pipeline. ControlNet에서 사용하는 RAFT 기반 "
        "optical flow warping을 frequency-aware temporal module로 교체 — "
        "optical flow 의존성을 완전히 제거할 수 있습니다.\n\n"
        "둘째 — DGAF-VSR과의 결합. 본 연구의 frequency-domain conditioning과 "
        "DGAF의 feature-domain warping은 orthogonal한 두 축이므로, "
        "fusion 시 cross-dataset 성능 개선이 기대됩니다.\n\n"
        "셋째 — 광범위한 training mixture. REDS와 Vimeo-90K 혼합 학습, "
        "그리고 degradation-aware augmentation으로 OOD gap을 완화할 수 "
        "있습니다.\n\n"
        "넷째 — latent-space wavelet analysis. Wavelet loss를 pixel "
        "space가 아닌 latent space에서 수행하고 adaptive band partitioning "
        "을 적용하면 VAE-decoding overhead를 크게 줄일 수 있습니다.\n\n"
        "다섯째 — short-sequence temporal regularization. 가벼운 flow-guided "
        "latent warping이나 multi-frame attention 같은 module을 추가해 "
        "UDM10에서 관찰된 flicker 문제를 해결할 수 있습니다.\n\n"
        "이상으로 발표를 마칩니다. 경청해 주셔서 감사합니다. 질문 부탁드립니다."
    ))


# ══════════════════════════════════════════════════════════════════════════
# Notes for the existing template slides (post-reorder, 0-indexed)
# ══════════════════════════════════════════════════════════════════════════
NOTES_EXISTING = {
    # 0-indexed positions for template slides in the final order
    9: (  # Architecture diagram (template slide 2)
        "이 슬라이드는 제안 framework의 전체 architecture를 보여줍니다.\n\n"
        "Panel (a) — inference pipeline 전체 입니다. LR 입력이 VAE encoder "
        "를 거쳐 latent space로 들어가고, U-Net이 diffusion denoising을 "
        "수행합니다. Temporal scaffold는 StableVSR를 그대로 계승했습니다 — "
        "ControlNet이 이전에 denoise된 이웃 프레임의 x̂_0를 RAFT optical "
        "flow로 warping해서 받습니다.\n\n"
        "새로 추가된 부분은 Wavelet Conditioning Module — WCM 입니다. "
        "WCM은 DT-CWT decomposition, Frequency Encoder, 그리고 BD-SFT "
        "injection으로 구성됩니다. 파란색으로 표시된 HIGH band modulation은 "
        "up_blocks[1] 즉 U-Net decoder에 들어가고, 빨간색으로 표시된 "
        "LOW band modulation은 down_blocks[1] 즉 encoder에 들어갑니다.\n\n"
        "Panel (b) — Frequency Encoder의 내부 입니다. 4개의 SubbandBlock "
        "이 있는데, 각각 DT-CWT scale L=1부터 4까지에 대응합니다. 다음 "
        "슬라이드에서 SubbandBlock 내부 detail을 다룹니다."
    ),
    12: (  # SubbandBlock diagram (template slide 3)
        "SubbandBlock 내부 — 논문에서는 PerLevelProcessor라고도 부릅니다.\n\n"
        "하나의 DT-CWT scale에서 추출한 6개 방향의 complex coefficient를 "
        "입력으로 받습니다. 두 개의 parallel branch — Magnitude Encoding과 "
        "Phase Encoding이 동일한 operator 시퀀스 (DWConv → SiLU → Spatial "
        "Attention → Conv-1x1)를 거쳐 처리됩니다.\n\n"
        "Recombination 단계에서 magnitude branch의 출력을 phase encoding "
        "의 cos과 sin과 element-wise 곱해 complex 임베딩의 real부와 "
        "imaginary부를 만듭니다. 이걸 magnitude와 함께 concatenate해 "
        "192 채널 feature가 되고, 두 개의 SFT head가 (γ, β) modulation "
        "parameter를 출력합니다.\n\n"
        "Identity-preserving initialization 덕분에 학습 시작 시점에는 "
        "wrapped U-Net이 원래 U-Net과 동일하게 동작하고, 학습이 진행되며 "
        "frequency-domain modulation을 점진적으로 학습합니다."
    ),
    21: (  # Evaluation visualization overview (template slide 4)
        "이 슬라이드는 본 연구의 qualitative 결과 종합 비교 입니다.\n\n"
        "왼쪽부터 Ours Evaluation, Ablation Study, Comparison, Frequency "
        "Average 영역으로 구성되어 있습니다. 본 방법은 StableVSR baseline "
        "보다 sharper texture와 cleaner edge를 보입니다 — 특히 repeated "
        "pattern이나 high-contrast structure 영역에서 frequency-domain "
        "conditioning이 spatial-domain ControlNet branch를 보완하는 효과가 "
        "확인됩니다."
    ),
    22: (  # LR subband visualization (template slide 5)
        "이 슬라이드는 LR 입력 이미지에 대한 DT-CWT subband 시각화 "
        "입니다.\n\n"
        "각 scale에서 ±15°, ±45°, ±75° 의 6개 directional subband의 "
        "magnitude를 보여줍니다. DT-CWT의 directional selectivity가 "
        "시각적으로 확인됩니다 — 각 subband는 해당 방향의 edge 에만 "
        "강하게 반응합니다. 이 directional richness가 decoder에서 "
        "directionally coherent한 HF detail 합성을 가능하게 합니다."
    ),
    23: (  # GT subband visualization (template slide 6)
        "이 슬라이드는 ground-truth HR 이미지에 대한 DT-CWT subband "
        "시각화 입니다.\n\n"
        "LR 대비 HR subband는 high-frequency content가 더 풍부하고 "
        "directional structure가 더 명확합니다. 본 연구의 wavelet loss는 "
        "예측된 subband를 이 GT subband와 매칭시키도록 supervision을 "
        "제공합니다."
    ),
    24: (  # Radial spectrum (template slide 7)
        "REDS4에서의 frequency-domain 분석 결과 입니다. Ground truth 대비 "
        "각 방법의 radial power spectrum 비율을 측정했습니다.\n\n"
        "세 가지 발견이 있습니다. 첫째 — 저주파에서는 모든 방법이 GT와 "
        "거의 일치 합니다. 둘째 — 중간 주파수에서는 본 방법이 비율 1에 "
        "거의 근접하지만 BasicVSR++와 StableVSR는 0.64 정도로 떨어집니다. "
        "셋째 — 고주파 영역에서 차이가 가장 큽니다. 본 방법은 GT 대비 "
        "0.69의 power retention을 보이지만 StableVSR는 0.34에 불과 — "
        "두 배 이상의 차이 입니다.\n\n"
        "이 frequency-domain 분석이 본 연구의 핵심 가설을 직접적으로 "
        "검증합니다 — HIGH band DT-CWT coefficient를 decoder로, LOW "
        "band를 encoder로 routing하는 설계가 pixel-level loss가 잡지 "
        "못하는 spectral content를 효과적으로 보존합니다."
    ),
    30: (  # Thanks
        "이상으로 발표를 마칩니다. 경청해 주셔서 감사합니다.\n\n"
        "질문 받겠습니다."
    ),
}


def reorder_slides(prs, new_order):
    sldIdLst = prs.slides._sldIdLst
    children = list(sldIdLst)
    assert sorted(new_order) == list(range(len(children)))
    for child in children:
        sldIdLst.remove(child)
    for idx in new_order:
        sldIdLst.append(children[idx])


def main():
    prs = Presentation(str(TEMPLATE))

    update_cover(prs)
    update_contents(prs)

    # New content slides — appended in build order (re-ordered below).
    # idx in slides after appending: 9 ..
    slide_motivation(prs)                # 9
    slide_related_work_paradigms(prs)    # 10
    slide_related_work_gap(prs)          # 11
    slide_contributions(prs)             # 12
    slide_prelim_vsr_diffusion(prs)      # 13
    slide_prelim_dwt_dtcwt(prs)          # 14
    slide_prelim_sft(prs)                # 15
    slide_why_dtcwt(prs)                 # 16
    slide_band_partition(prs)            # 17
    slide_subbandblock_detail(prs)       # 18
    slide_bdsft_injection(prs)           # 19
    slide_training_loss(prs)             # 20
    slide_experiments_setup(prs)         # 21
    slide_metrics(prs)                   # 22
    slide_results_reds4(prs)             # 23
    slide_results_dm_group(prs)          # 24
    slide_results_cross_dataset(prs)     # 25
    slide_ablation(prs)                  # 26
    slide_discussion(prs)                # 27
    slide_limitations(prs)               # 28
    slide_conclusion(prs)                # 29
    slide_future_work(prs)               # 30

    # Existing template slides indices: 0 Cover, 1 Contents, 2 Arch,
    # 3 SubbandBlock, 4 Eval overview, 5 LR subbands, 6 GT subbands,
    # 7 Radial spectrum, 8 Thanks.
    final_order = [
        0,    # 1. Cover
        1,    # 2. Contents
        9,    # 3. Motivation
        10,   # 4. Related Work — Paradigms
        11,   # 5. Related Work — Gap
        12,   # 6. Contributions
        13,   # 7. Preliminaries · VSR & Diffusion
        14,   # 8. Preliminaries · DWT → DT-CWT
        15,   # 9. Preliminaries · SFT
        2,    # 10. Architecture diagram (template)
        16,   # 11. Why DT-CWT
        17,   # 12. Band partitioning
        3,    # 13. SubbandBlock diagram (template)
        18,   # 14. SubbandBlock detail
        19,   # 15. BD-SFT Injection
        20,   # 16. Training Loss
        21,   # 17. Datasets & Setup
        22,   # 18. Metrics
        23,   # 19. REDS4 vs StableVSR
        24,   # 20. DM Group
        25,   # 21. Cross-Dataset
        4,    # 22. Eval overview (template)
        5,    # 23. LR subbands (template)
        6,    # 24. GT subbands (template)
        7,    # 25. Radial spectrum (template)
        26,   # 26. Ablation
        27,   # 27. Discussion
        28,   # 28. Limitations
        29,   # 29. Conclusion
        30,   # 30. Future Work
        8,    # 31. Thanks
    ]
    reorder_slides(prs, final_order)

    # Attach notes to existing template slides at their final positions.
    for final_idx, note in NOTES_EXISTING.items():
        if final_idx < len(prs.slides):
            set_notes(prs.slides[final_idx], note)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
