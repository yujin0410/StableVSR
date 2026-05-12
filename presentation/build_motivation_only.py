"""Build a single-slide PPT containing ONLY the redesigned Motivation slide.

User can copy-paste this slide into the main deck. The file uses the same
template, layout 12 ('16_제목만'), and stylistic conventions as the main
deck (no top section badges; section accent = blue Intro color).
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


TEMPLATE = Path("/tmp/thesis_work/template.pptx")
FIG_DIR = Path("/tmp/thesis_work/fig_png")
OUT = Path("/home/user/StableVSR/presentation/motivation_slide_only.pptx")


# ── Palette (mirrors the main deck) ────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x12, 0x1E, 0x3D)
MUTED_INK = RGBColor(0x3D, 0x3D, 0x4A)
GRAY = RGBColor(0x6B, 0x6B, 0x78)
LIGHT_GRAY = RGBColor(0xE5, 0xE5, 0xE5)
SOFT_BG2 = RGBColor(0xF9, 0xF9, 0xFC)
ACCENT = RGBColor(0x2C, 0x55, 0xA8)   # Section 1 — Intro blue


def _rm(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def add_card(slide, left, top, width, height, *, fill=WHITE, line=None,
              line_width=0.5):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(top),
                                Inches(width), Inches(height))
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_width)
    s.text_frame.text = ""
    return s


def add_strip(slide, left, top, width, color, *, thickness=0.045):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(left), Inches(top),
                                Inches(width), Inches(thickness))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False


def add_left_stripe(slide, left, top, height, color, *, thickness=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(left), Inches(top),
                                Inches(thickness), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False


def add_textbox(slide, left, top, width, height, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    return box


def add_para(tf, text, *, size=13, bold=False, italic=False, color=INK,
              align=PP_ALIGN.LEFT, space_before=2, space_after=4):
    if not tf.text and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    if color is not None:
        r.font.color.rgb = color
    return p


def stripe_card(slide, left, top, w, h, *, accent, fill=WHITE,
                 stripe_top=True, stripe_left=False):
    add_card(slide, left, top, w, h, fill=fill, line=LIGHT_GRAY)
    if stripe_top:
        add_strip(slide, left, top, w, accent, thickness=0.05)
    if stripe_left:
        add_left_stripe(slide, left, top, h, accent, thickness=0.07)


def set_title(slide, title, subtitle=""):
    ph = None
    for cand in slide.placeholders:
        if cand.placeholder_format.idx == 1:
            ph = cand
            break
    tf = ph.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    main = p.add_run()
    main.text = title
    main.font.size = Pt(22)
    main.font.bold = True
    main.font.color.rgb = WHITE
    if subtitle:
        sub = p.add_run()
        sub.text = f"   ·   {subtitle}"
        sub.font.size = Pt(14)
        sub.font.color.rgb = RGBColor(0xDD, 0xD8, 0xF0)


def add_picture_fit(slide, path, left, top, width, height):
    from PIL import Image
    with Image.open(path) as im:
        w_px, h_px = im.size
    ai = w_px / h_px
    ab = width / height
    if ai > ab:
        new_w = width
        new_h = width / ai
    else:
        new_h = height
        new_w = height * ai
    L = left + (width - new_w) / 2
    T = top + (height - new_h) / 2
    slide.shapes.add_picture(str(path), Inches(L), Inches(T),
                              Inches(new_w), Inches(new_h))


def main():
    prs = Presentation(str(TEMPLATE))

    # Drop the template's pre-existing slides; we only need a blank canvas
    # using layout 12 ('16_제목만').
    sldIdLst = prs.slides._sldIdLst
    for child in list(sldIdLst):
        sldIdLst.remove(child)

    s = prs.slides.add_slide(prs.slide_layouts[12])
    # Drop layout's default body placeholder
    for ph in list(s.placeholders):
        if ph.placeholder_format.idx == 10:
            _rm(ph)

    set_title(s, "Motivation", "Why VSR + Diffusion, and why is it hard?")

    # ── Row 1 — VSR task (B option: diagram + bullets) ─────────────
    stripe_card(s, 0.55, 1.0, 12.3, 2.5, accent=ACCENT)
    # Diagram on the left ~7"
    add_picture_fit(s, FIG_DIR / "motivation_vsr_task.png",
                     0.7, 1.10, 7.7, 2.3)
    # Right-side bullet box
    tb = add_textbox(s, 8.7, 1.15, 4.0, 2.25)
    add_para(tb.text_frame, "Video Super-Resolution",
              size=13, bold=True, color=ACCENT)
    add_para(tb.text_frame, "Input  ·  LR frame sequence",
              size=11.5, bold=True, color=INK, space_before=6)
    add_para(tb.text_frame, "{ I^LR_{t−1},  I^LR_t,  I^LR_{t+1} }",
              size=11, color=MUTED_INK, space_before=1)
    add_para(tb.text_frame, "Output  ·  HR frame",
              size=11.5, bold=True, color=INK, space_before=6)
    add_para(tb.text_frame, "Î^HR_t",
              size=11, color=MUTED_INK, space_before=1)
    add_para(tb.text_frame, "Key  ·  temporal context across frames",
              size=11.5, bold=True, color=INK, space_before=6)
    add_para(tb.text_frame,
              "single-image SR과 달리 인접 프레임 정보를 함께 활용 — "
              "한 프레임에서 가려진 detail을 이웃이 보완.",
              size=10.5, italic=True, color=GRAY, space_before=1)

    # ── Row 2 — Flicker problem + diagram ──────────────────────────
    LOW = RGBColor(0xCB, 0x42, 0x35)
    stripe_card(s, 0.55, 3.70, 12.3, 2.45, accent=LOW)
    # Left summary text
    tb = add_textbox(s, 0.75, 3.85, 4.0, 2.25)
    add_para(tb.text_frame, "But: temporal flickering",
              size=13, bold=True, color=LOW)
    add_para(tb.text_frame,
              "Diffusion prior로 사실적인 texture를 합성하지만, "
              "프레임 독립 stochastic denoising 때문에 HF detail이 "
              "frame 간 불일치 → flicker.",
              size=11, color=MUTED_INK, space_before=8)
    add_para(tb.text_frame, "본 연구가 해결하려는 핵심 bottleneck.",
              size=11, italic=True, color=LOW, space_before=8)
    # Flicker diagram on the right
    add_picture_fit(s, FIG_DIR / "motivation_flicker.png",
                     4.9, 3.85, 7.9, 2.25)

    # ── Row 3 — Research question ─────────────────────────────────
    add_card(s, 0.55, 6.30, 12.3, 0.95, fill=SOFT_BG2, line=LIGHT_GRAY)
    add_left_stripe(s, 0.55, 6.30, 0.95, ACCENT, thickness=0.10)
    tb = add_textbox(s, 0.85, 6.40, 11.85, 0.85, anchor=MSO_ANCHOR.MIDDLE)
    add_para(tb.text_frame, "Research question",
              size=11, bold=True, color=ACCENT)
    add_para(tb.text_frame,
              "사전학습된 image diffusion model을 parameter-efficient하게 "
              "video에 adapting하면서도, 3D attention · full fine-tuning 없이 "
              "temporal consistency를 달성할 수 있을까?",
              size=12, bold=True, color=INK, space_before=2)

    # ── Speaker notes ──────────────────────────────────────────────
    notes_tf = s.notes_slide.notes_text_frame
    notes_tf.clear()
    p = notes_tf.paragraphs[0]
    p.add_run().text = (
        "Motivation 슬라이드입니다. 세 구역으로 구성되어 있습니다.\n\n"
        "첫째, VSR이라는 task 자체를 설명합니다. 그림에서 보시는 것처럼, "
        "VSR은 저해상도 프레임 시퀀스 — t-1, t, t+1 — 을 입력으로 받아 "
        "고해상도 프레임 한 장을 출력합니다. Single-image SR과의 결정적 "
        "차이는 인접 프레임의 temporal context를 함께 활용한다는 점 "
        "입니다. 한 프레임에서 흐리거나 가려진 detail을 다른 프레임이 "
        "보완할 수 있기 때문에 단일 영상 SR보다 본질적으로 더 풍부한 "
        "정보를 활용합니다.\n\n"
        "둘째, 그러나 diffusion prior를 video에 그대로 가져오면 문제가 "
        "생깁니다. 프레임마다 독립적으로 stochastic denoising을 수행하기 "
        "때문에, 같은 영역의 high-frequency detail이 프레임 간에 일관되지 "
        "않습니다. 슬라이드 오른쪽 그림 — 윗줄을 보면 같은 장면임에도 "
        "프레임마다 texture가 미세하게 다르게 합성되어 어른거리는 flicker "
        "가 발생합니다. 본 연구의 목표는 아랫줄처럼 frame 간 detail이 "
        "안정적으로 유지되도록 만드는 것입니다.\n\n"
        "셋째, 본 연구의 research question은 — 사전학습된 image diffusion "
        "model을 parameter-efficient하게 video에 적응시키면서도, 3D "
        "attention이나 full-network fine-tuning 없이 temporal consistency "
        "를 달성할 수 있을까 — 입니다. 답은 'pixel을 warping하지 말고, "
        "shift-stable한 frequency-domain prior로 conditioning하자' 입니다. "
        "다음 슬라이드부터 자세히 설명드리겠습니다."
    )

    prs.save(str(OUT))
    print(f"Saved single-slide PPT: {OUT}")


if __name__ == "__main__":
    main()
