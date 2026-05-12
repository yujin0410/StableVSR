"""Translate all Korean text in slide BODIES to English.

  - Speaker notes (slide.notes_slide) are NOT touched.
  - Matches run.text exactly against a hand-curated translation map.
  - Any Korean-containing run not in the map is reported so we can
    iterate until the file is fully bilingual.
"""

from pathlib import Path
import sys

from pptx import Presentation


SRC = Path("/tmp/user_combined.pptx")
DST = Path("/home/user/StableVSR/presentation/WC_BD_SFT_Defense_Cho_YuJin_EN.pptx")


def has_korean(s: str) -> bool:
    return any('가' <= ch <= '힣' for ch in s)


# ── Translation table (run.text → English) ─────────────────────────
# Keys must match the run text EXACTLY (including punctuation, leading
# bullets like '▸  ', and embedded \n).
TRANS = {
    # Slide 1 — cover
    "석사학위논문 예비심사":
        "Master's Thesis Preliminary Defense",
    "공학과 조유진":
        " Engineering, Yu-Jin Cho",

    # Slide 3 — motivation (3 cards + research question)
    "고품질 HR 영상 수요 급증 · VSR이 spatial · temporal 정보로 HR 복원.":
        "Demand for HR video surging · VSR recovers HR via spatial · temporal info.",
    "CNN/GAN의 over-smoothing · instability를 넘어 realistic texture (SR3, DDPM).":
        "Beyond CNN/GAN over-smoothing · instability — realistic texture (SR3, DDPM).",
    "프레임 독립 stochastic denoising ⇒ HF detail이 frame 간 불일치 → flicker.":
        "Per-frame independent stochastic denoising ⇒ HF detail inconsistent across frames → flicker.",
    "사전학습된 image diffusion model을 parameter-efficient하게 video에 adapting하면서도, 3D attention · full fine-tuning 없이 temporal consistency를 달성할 수 있을까?":
        "Can we adapt a pre-trained image diffusion model to video parameter-efficiently, achieving temporal consistency without 3D attention or full fine-tuning?",

    # Slide 4 — motivation v2 (option B layout)
    "single-image SR과 달리 인접 프레임 정보를 함께 활용 — 한 프레임에서 가려진 detail을 이웃이 보완.":
        "Unlike single-image SR, uses adjacent-frame info — detail occluded in one frame is recovered from neighbours.",
    "Diffusion prior로 사실적인 texture를 합성하지만, 프레임 독립 stochastic denoising 때문에 HF detail이 frame 간 불일치 → flicker.":
        "Diffusion prior synthesizes realistic texture, but per-frame independent stochastic denoising leaves HF detail inconsistent across frames → flicker.",
    "본 연구가 해결하려는 핵심 bottleneck.":
        "The core bottleneck this work addresses.",

    # Slide 6 — Related Work · Gap
    "Video temporal consistency 문제를 다루지 않음. Diffusion 기반 frequency 연구들은 모두 image SR만 대상.":
        "Does not address video temporal consistency. Existing diffusion-based frequency works target image SR only.",
    "Critical decimation 때문에 sub-pixel motion에도 wavelet coefficient가 크게 흔들림 → video flickering 악화.":
        "Critical decimation makes wavelet coefficients fluctuate even under sub-pixel motion → worsens video flickering.",
    "Image SR이 아닌 video VSR을 대상으로, DWT가 아닌 near-shift-invariant DT-CWT를 사용해, ":
        "Target video VSR rather than image SR; use near-shift-invariant DT-CWT rather than DWT; ",
    "사전학습된 diffusion prior를 freeze한 채 asymmetric encoder-decoder injection으로 conditioning한다.":
        "Freeze the pre-trained diffusion prior and condition via asymmetric encoder-decoder injection.",

    # Slide 7 — Contributions
    "▸  DT-CWT 기반 frequency-domain conditioning을 frozen pre-trained diffusion U-Net에 주입.":
        "▸  Inject DT-CWT-based frequency-domain conditioning into a frozen pre-trained diffusion U-Net.",
    "▸  Trigonometric phase encoding — 2π wraparound 불연속 회피.":
        "▸  Trigonometric phase encoding — avoids the 2π-wraparound discontinuity.",
    "▸  HIGH는 magnitude only, LOW는 magnitude-weighted angular distance.":
        "▸  HIGH: magnitude only.  LOW: magnitude-weighted angular distance.",
    "▸  Latent ε-MSE에 더해 pixel-space frequency supervision.":
        "▸  Pixel-space frequency supervision on top of latent ε-MSE.",

    # Slide 8 — Preliminaries · VSR & Diffusion
    "→  본 연구는 이 scaffold를 그대로 두고 WCM만 추가.":
        "→  This work keeps the scaffold intact and adds only the WCM.",
    "Blau & Michaeli (2018) — pixel fidelity와 perceptual quality는 동시 최대화 불가.":
        "Blau & Michaeli (2018) — pixel fidelity and perceptual quality cannot be maximised together.",
    "본 연구는 perceptual 극단을 추구하는 설계.":
        "This work targets the perceptual extreme by design.",

    # Slide 9 — DWT → DT-CWT
    "✗  Shift-variance — critical decimation 때문에 sub-pixel shift에도 coefficient가 크게 흔들림.":
        "✗  Shift-variance — critical decimation makes coefficients fluctuate even under sub-pixel shifts.",
    "✗  3 directional subbands per scale — horizontal · vertical · diagonal만.":
        "✗  Only 3 directional subbands per scale — horizontal · vertical · diagonal.",
    "✗  Real-valued — phase 정보 미포함.":
        "✗  Real-valued — no phase information.",
    "✓  Near shift-invariance — 두 parallel tree (Hilbert transform pair).":
        "✓  Near shift-invariance — two parallel trees (Hilbert transform pair).",
    "✓  Complex coefficients C = a + ib  ⇒  (magnitude, phase) 분해:":
        "✓  Complex coefficients C = a + ib  ⇒  decompose into (magnitude, phase):",

    # Slide 10 — SFT preliminaries
    "feature를 channel-wise 가 아닌 spatial-wise로 affine modulation.":
        "Spatially-varying (not channel-wise) affine modulation of features.",
    "Globally uniform.\nSpatial detail 손실.":
        "Globally uniform.\nLoses spatial detail.",
    "Spatially varying + linear-time.\nWavelet conditioning에 최적.":
        "Spatially varying + linear-time.\nIdeal for wavelet conditioning.",
    "▸  Wavelet coefficient는 본질적으로 spatial하게 분포 — spatial-wise modulation이 자연스러운 매칭.":
        "▸  Wavelet coefficients are inherently spatial — spatial-wise modulation is a natural fit.",
    "▸  Cross-attention 대비 가벼움 — H × W 크기의 latent에서 quadratic cost 회피.":
        "▸  Lighter than cross-attention — avoids quadratic cost on H × W latents.",
    "▸  Identity-preserving init (γ=1, β=0)이 자연스럽게 정의됨 — 학습 시작 시 baseline 동작 보존.":
        "▸  Identity-preserving init (γ=1, β=0) is naturally defined — preserves baseline behaviour at step 0.",

    # Slide 12 — Why DT-CWT
    "DWT는 critical decimation으로 frame 간 inconsistent conditioning.":
        "DWT, via critical decimation, produces inconsistent conditioning across frames.",
    "Hair / foliage / fabric의 임의 방향 edge 보존 — oriented HF detail 합성에 적합.":
        "Preserves arbitrary-orientation edges in hair / foliage / fabric — ideal for oriented HF detail synthesis.",
    "φ = ∠C → structure (precise spatial localization) — BD-SFT pathway에 직접 대응.":
        "φ = ∠C → structure (precise spatial localization) — directly maps to the BD-SFT pathway.",

    # Slide 13 — Band partition rationale
    "▸  Spectral — LR 신호 spectrum을 정확히 절반으로 분리 (texture 상위, structure 하위). 각 SFT pathway가 non-overlapping range에 특화.":
        "▸  Spectral — exactly halves the LR spectrum (texture upper, structure lower); each SFT pathway specialises on a non-overlapping range.",
    "▸  Architectural — 2 scales × 256 ch = 512 ch ⇒ U-Net의 block_out_channels[1] = 512와 정확히 일치, 추가 projection 불필요.":
        "▸  Architectural — 2 scales × 256 ch = 512 ch matches U-Net's block_out_channels[1] = 512 exactly; no extra projection needed.",
    "▸  J = 4 — j=4 subband의 H/16 × W/16 size가 SD-VAE 8× 와 U-Net 2× 다운샘플링 위에서 자연스럽게 정렬.":
        "▸  J = 4 — the j=4 subband's H/16 × W/16 size aligns naturally on top of SD-VAE 8× and U-Net 2× downsampling.",

    # Slide 15 — SubbandBlock detail
    "▸  ±π 경계에서 continuous · gradient stable":
        "▸  Continuous across ±π · gradient-stable",
    "▸  ε = 1e-8 — flat region에서 M → 0 안정화":
        "▸  ε = 1e-8 — stabilises M → 0 in flat regions",

    # Slide 17 — Training Loss
    "Texture를 loosely supervise — network가 phase에 constraint 없이 자유롭게 detail을 합성하도록.":
        "Loosely supervise texture — network is free to synthesise detail without phase constraints.",
    "Structural phase 오차를 |C| 비례로 강력하게 supervise — high-energy 영역의 phase가 더 중요.":
        "Strongly supervise structural phase error weighted by |C| — phase matters more in high-energy regions.",

    # Slide 18 — Setup
    "▸  유일한 architectural 차이 — Wavelet Conditioning Module (WCM) + BD-SFT injection.":
        "▸  Only architectural difference — Wavelet Conditioning Module (WCM) + BD-SFT injection.",
    "⇒  REDS4 performance gap이 WCM에 의해 isolation됨 — ablation-by-design.":
        "⇒  REDS4 performance gap is isolated to the WCM — ablation-by-design.",

    # Slide 19 — Metrics
    "Generative VSR은 perception–distortion trade-off 위에 위치 — 단일 metric으로 전체 그림을 보여줄 수 없습니다. Pixel metric은 사실적인 합성을 페널라이즈하고, NR perceptual metric은 reference-based metric을 보완하며, temporal metric은 per-frame metric이 놓치는 flicker를 포착합니다.":
        "Generative VSR sits on the perception–distortion trade-off — no single metric captures the full picture. Pixel metrics penalise realistic synthesis; NR perceptual metrics complement reference-based ones; temporal metrics expose flicker that per-frame metrics miss.",

    # Slide 20 — REDS4
    "Training data · backbone · temporal scaffold이 모두 matched된 상태이므로, 이 성능 gap은 wavelet-conditioned BD-SFT mechanism의 단독 기여로 정확하게 attribute됩니다.":
        "Training data · backbone · temporal scaffold are all matched, so this performance gap attributes directly to the wavelet-conditioned BD-SFT mechanism.",

    # Slide 21 — DM Group
    "PSNR / SSIM은 BasicVSR++가 우수 (regression).":
        "BasicVSR++ leads on PSNR / SSIM (regression target).",
    "Perceptual 4 metric에서는 모두 열위 — Blau & Michaeli (2018)의 perception-distortion trade-off.":
        "Trails on all 4 perceptual metrics — Blau & Michaeli (2018) perception-distortion trade-off.",
    "DM 내부 ranking이 본 연구 mechanism의 공정한 비교 축.":
        "Within-DM ranking is the fair comparison axis for our mechanism.",

    # Slide 22 — Cross-dataset
    "tLPIPS 51.11 → 31.59 (−38.2 %) over StableVSR.\nLong-motion content에서는 frequency conditioning이 stability에 도움.":
        "tLPIPS 51.11 → 31.59 (−38.2 %) over StableVSR.\nFor long-motion content, frequency conditioning helps stability.",
    "Compressed SD vs. clean REDS 720p — gain 감소.\nDGAF의 OGWM과 fusion이 future work.":
        "Compressed SD vs. clean REDS 720p — gains shrink.\nFusion with DGAF's OGWM is future work.",

    # Slide 28 — Discussion
    "DM 방법 → high perceptual quality.":
        "DM methods → high perceptual quality.",
    "WC-BD-SFT는 perceptual 극단 — REDS4 / UDM10 NR 1위.":
        "WC-BD-SFT at the perceptual extreme — #1 NR on REDS4 / UDM10.",
    "REDS-only 학습 → WCM 기여 isolation.":
        "REDS-only training → isolates WCM contribution.",
    "REDS-like statistics에 강함, Vid4 SD에서는 gain 감소.":
        "Strong on REDS-like statistics; smaller gains on Vid4 SD.",
    "UDM10 (32-frame, 정적) → tLPIPS 악화.":
        "UDM10 (32-frame, static) → tLPIPS degrades.",
    "SPMCS (long-motion) → tLPIPS 개선.":
        "SPMCS (long-motion) → tLPIPS improves.",

    # Slide 29 — Limitations
    "REDS-like content에서 최강 · Vid4 SD에서 gain 감소.":
        "Strongest on REDS-like content; smaller gains on Vid4 SD.",
    "→ 광범위한 학습 mix · degradation-aware augmentation":
        "→ Broader training mixture · degradation-aware augmentation",
    "PSNR/SSIM을 perceptual quality에 일부 양보 — 특히 UDM10.":
        "PSNR/SSIM partly yielded to perceptual quality — especially on UDM10.",
    "→ Perceptual application에 적합 · 의료/과학 영상엔 부적합":
        "→ Suited for perceptual applications · not for medical/scientific imaging",
    "UDM10 (32 frame)에서 DGAF/StableVSR보다 tLPIPS 높음.":
        "UDM10 (32 frame) — tLPIPS higher than DGAF/StableVSR.",
    "4-level DT-CWT가 frame당 modest cost 추가.":
        "4-level DT-CWT adds modest per-frame cost.",
    "→ Diffusion denoising cost에 비해 작음 · 추가 최적화 가능":
        "→ Small vs. diffusion denoising cost · further optimisation possible",
    "Single fixed seed · 2025년 다수 method 미평가.":
        "Single fixed seed · many 2025 methods not evaluated.",
    "→ Multi-seed evaluation · RVRT / MGLD / DiffVSR 후속 비교":
        "→ Multi-seed evaluation · follow-up RVRT / MGLD / DiffVSR comparison",

    # Slide 31 — Future Work
    "ControlNet의 RAFT warping을 frequency-aware temporal module로 교체 — optical-flow 의존성 완전 제거.":
        "Replace ControlNet's RAFT warping with a frequency-aware temporal module — fully removes optical-flow dependency.",
    "Frequency-domain conditioning ⟂ feature-domain warping — 두 방향이 orthogonal해 fusion 시 cross-dataset 개선 기대.":
        "Frequency-domain conditioning ⟂ feature-domain warping — orthogonal axes; fusion expected to improve cross-dataset performance.",
    "REDS + Vimeo-90K 학습 · degradation-aware augmentation — Vid4 같은 OOD에서의 gap 완화.":
        "Train on REDS + Vimeo-90K · degradation-aware augmentation — narrows the gap on OOD content like Vid4.",
    "Wavelet loss를 latent space에서 수행 · adaptive band partitioning — VAE-decoding overhead 감소.":
        "Apply the wavelet loss in latent space · adaptive band partitioning — reduces VAE-decoding overhead.",
    "Flow-guided latent warping / multi-frame attention 같은 경량 module 추가 — UDM10 flicker 해결.":
        "Add lightweight modules (flow-guided latent warping / multi-frame attention) — addresses UDM10 flicker.",

    # Slide 32 — References footnote
    "전체 인용 목록은 학위논문의 reference.bib 참조 (73 entries).":
        "Full citation list: see reference.bib in the thesis (73 entries).",
}


def translate_slides(src_path: Path, dst_path: Path):
    prs = Presentation(str(src_path))
    misses = []
    hits = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        # ONLY touch shapes; do not enter slide.notes_slide.
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if not has_korean(run.text):
                        continue
                    key = run.text
                    if key in TRANS:
                        run.text = TRANS[key]
                        hits += 1
                    else:
                        misses.append((slide_idx, key))

    prs.save(str(dst_path))
    print(f"Translated {hits} run(s).")
    if misses:
        print(f"\nMISSED ({len(misses)}) — add these to TRANS map:")
        for s_i, t in misses:
            print(f"  [S{s_i}] {repr(t)}")
    else:
        print("All Korean runs in slide bodies translated.")


if __name__ == "__main__":
    translate_slides(SRC, DST)
